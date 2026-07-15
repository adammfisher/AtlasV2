"""Layout-archetype exemplar engine (v2).

dfs_exemplars.json carries 1–2 schema-valid exemplar SPECS per archetype —
assertive sentence headlines, parallel bullet grammar, quantified evidence,
speaker notes — each anchored to a real slide in dfs_library.pptx via
geometry_source (extracted geometry, never invented). retrieve_exemplars()
selects the top-K whose tags match a request's content shape; the server
injects them into the generation prompt (small tiers lean on them hardest —
k defaults low to respect small-model context budgets).

The TypeScript mirror (server/src/pipeline/exemplars.ts) implements the same
scoring over the same manifest for prompt assembly in the app server (which
ships no Python); this module is canonical for tests, evals, and tooling.
"""
import json
import re
from pathlib import Path

HERE = Path(__file__).resolve().parents[2] / "skills/pptx/templates"
_BUNDLED = Path(__file__).resolve().parent / "templates"

_STOP = frozenset(
    "a an and are as at be by for from has have how in is it of on or our the "
    "this that to was we what which with your make makes made get give show me my".split()
)

# content-shape hints: request phrasing → archetype affinity
_SHAPE_HINTS = {
    "content_chart": ["chart", "graph", "trend", "over time", "monthly", "quarterly", "growth", "revenue", "funnel", "metrics"],
    "big_stat": ["metric", "kpi", "number", "record", "milestone", "headline", "roi", "savings"],
    "comparison": ["versus", "vs", "compare", "comparison", "competitor", "options", "pros", "cons", "before", "after"],
    "timeline_process": ["timeline", "roadmap", "phases", "steps", "plan", "rollout", "launch", "process", "sequence"],
    "table": ["table", "dashboard", "exact", "targets", "actuals", "breakdown"],
    "quote": ["quote", "testimonial", "customer said", "voice", "feedback"],
    "two_column": ["screenshot", "feature", "side by side", "narrative"],
    "content_bullets": ["points", "reasons", "findings", "drivers", "risks", "summary"],
    "section_divider": ["sections", "parts", "chapters"],
    "agenda": ["agenda", "overview"],
    "title": ["deck", "presentation", "review", "pitch"],
    "closing_cta": ["ask", "decision", "next steps", "closing"],
}


def _load_manifest() -> dict:
    for base in (HERE, _BUNDLED):
        path = base / "dfs_exemplars.json"
        if path.exists():
            return json.loads(path.read_text())
    raise FileNotFoundError("dfs_exemplars.json not found")


def _tokens(text: str) -> set:
    return {t for t in re.findall(r"[a-z0-9%$]+", str(text).lower()) if t not in _STOP}


def retrieve_exemplars(spec_request, k: int = 3) -> list:
    """Top-K exemplars for a request. `spec_request` is the user's request text
    (or a dict with a "text" key). Scoring is deterministic: tag-token overlap
    + content-shape hint hits, with archetype diversity (max one exemplar per
    archetype until every scored archetype is represented)."""
    text = spec_request.get("text") if isinstance(spec_request, dict) else spec_request
    req_tokens = _tokens(text or "")
    low = str(text or "").lower()
    manifest = _load_manifest()

    scored = []
    for exemplar in manifest["exemplars"]:
        tag_hits = sum(1 for tag in exemplar["tags"] if _tokens(tag) & req_tokens)
        hint_hits = sum(1 for hint in _SHAPE_HINTS.get(exemplar["archetype"], []) if hint in low)
        score = tag_hits * 2 + hint_hits
        scored.append((score, exemplar["id"], exemplar))
    scored.sort(key=lambda t: (-t[0], t[1]))

    picked, seen_archetypes = [], set()
    for score, _, exemplar in scored:  # diversity pass: one per archetype first
        if score > 0 and exemplar["archetype"] not in seen_archetypes:
            picked.append(exemplar)
            seen_archetypes.add(exemplar["archetype"])
        if len(picked) == k:
            return picked
    for score, _, exemplar in scored:  # backfill: next-best regardless of archetype
        if exemplar not in picked:
            picked.append(exemplar)
        if len(picked) == k:
            break
    return picked


def format_exemplars(exemplars: list) -> str:
    """Prompt block: compact spec JSON with the why, per exemplar."""
    parts = []
    for e in exemplars:
        parts.append(
            f"### {e['archetype']} — {e['why_good']}\n{json.dumps(e['spec'], separators=(',', ':'))}"
        )
    return "\n".join(parts)


# ── dev utility: regenerate template_geometry from the real library ─────────
def extract_geometry(library_path: Path, anchors: dict) -> dict:
    """Extract headline/content/graphics rects per archetype from anchor slides
    in dfs_library.pptx — the manifest's geometry is measured, never invented."""
    from pptx import Presentation

    emu = 914400.0
    prs = Presentation(str(library_path))
    out = {}
    for archetype, idx in anchors.items():
        infos = []

        def walk(shapes):
            for sh in shapes:
                if sh.shape_type == 6 and hasattr(sh, "shapes"):
                    walk(sh.shapes)
                    continue
                try:
                    rect = [round(sh.left / emu, 2), round(sh.top / emu, 2),
                            round(sh.width / emu, 2), round(sh.height / emu, 2)]
                except Exception:
                    continue
                if getattr(sh, "has_chart", False):
                    infos.append({"kind": "chart", "rect": rect})
                elif getattr(sh, "has_table", False):
                    infos.append({"kind": "table", "rect": rect})
                elif getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                    sizes = [r.font.size.pt for p in sh.text_frame.paragraphs
                             for r in p.runs if r.font.size]
                    infos.append({"kind": "text", "rect": rect,
                                  "max_pt": max(sizes) if sizes else None,
                                  "sample": sh.text_frame.text.strip().replace("\n", " | ")[:60]})

        walk(prs.slides[idx].shapes)
        texts = [i for i in infos if i["kind"] == "text" and i["rect"][1] < 6.5]
        texts.sort(key=lambda i: ((i["max_pt"] or 0), -i["rect"][1]), reverse=True)
        out[archetype] = {
            "library_slide": idx,
            "headline": texts[0] if texts else None,
            "content": texts[1:5],
            "graphics": [i for i in infos if i["kind"] in ("chart", "table")],
        }
    return out

"""
atlasv2-office Lambda (PRD §12.1): scale-to-zero Python office builder. The
main app Lambda invokes this with {skill, payload, template_b64?}; it runs the
same build_<skill>.py used locally and returns {ok, meta, checks, file_b64}.
Single-file kinds only (pptx/docx/xlsx/pdf); react/site/mermaid/svg/md never
touch Python. Pure-python deps in a layer — no containers, no native libs.
"""
import base64
import importlib
import json
import os
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

SKILLS = {"pptx", "docx", "xlsx", "pdf"}
# template filename per skill (bundled beside the handler under templates/)
TEMPLATES = {"pptx": "dfs_default.potx", "docx": "atlas_default.dotx"}


def _theme_colors(prs):
    """Resolve the deck's theme palette (dk1/lt1/accent1..6/…) → hex from the
    first master's theme part. Without this, theme-colored fills and text (which
    python-pptx can't give an .rgb for) vanish — e.g. a full-bleed accent title
    background renders as blank. Best-effort; returns {} on any surprise."""
    from lxml import etree

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    out = {}
    try:
        master = prs.slide_masters[0]
        # the theme is a generic Part (no python-pptx .element) — parse its raw XML
        theme_part = next((r.target_part for r in master.part.rels.values() if "theme" in r.reltype), None)
        if theme_part is None:
            return out
        scheme = etree.fromstring(theme_part.blob).find(".//a:clrScheme", ns)
        for child in (scheme if scheme is not None else []):
            tag = etree.QName(child).localname
            srgb = child.find("a:srgbClr", ns)
            sysc = child.find("a:sysClr", ns)
            if srgb is not None and srgb.get("val"):
                out[tag] = "#" + srgb.get("val")
            elif sysc is not None:
                out[tag] = "#" + (sysc.get("lastClr") or "000000")
    except Exception:
        pass
    return out


def _wrap(text, max_chars):
    """Greedy word-wrap to an estimated chars-per-line, so multi-line text frames
    render with the same line breaks PowerPoint would apply (the single biggest
    fidelity gap in the pure-SVG preview)."""
    words, lines, cur = text.split(), [], ""
    for wd in words:
        if not cur:
            cur = wd
        elif len(cur) + 1 + len(wd) <= max_chars:
            cur += " " + wd
        else:
            lines.append(cur)
            cur = wd
    if cur:
        lines.append(cur)
    return lines or [""]


def _apply_transforms(r, g, b, transforms):
    """Apply OOXML color transforms (lumMod/lumOff/tint/shade/satMod) so a themed
    accent used at, say, lumMod 20% + lumOff 80% renders as its true light tint
    instead of the full-saturation base — the difference between a light-pink
    card and a dark-magenta one."""
    if not transforms:
        return r, g, b
    import colorsys
    h, l, s = colorsys.rgb_to_hls(r / 255.0, g / 255.0, b / 255.0)
    for t, v in transforms:
        if t == "lumMod":
            l *= v
        elif t == "lumOff":
            l += v
        elif t == "satMod":
            s *= v
        elif t == "tint":  # lighten toward white
            l = l + (1.0 - l) * (1.0 - v)
        elif t == "shade":  # darken toward black
            l = l * v
    l = min(1.0, max(0.0, l))
    s = min(1.0, max(0.0, s))
    rr, gg, bb = colorsys.hls_to_rgb(h, l, s)
    return int(rr * 255), int(gg * 255), int(bb * 255)


def _clr_hex(el, theme, default=None):
    """Resolve an <a:srgbClr>/<a:schemeClr> element to hex, honoring the theme
    palette and any luminance/tint transforms on it."""
    if el is None:
        return default
    from lxml import etree

    tag = etree.QName(el).localname
    if tag == "srgbClr":
        base = el.get("val")
    elif tag == "schemeClr":
        val = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}.get(el.get("val"), el.get("val"))
        hx = theme.get(val)
        base = hx[1:] if hx else None
    else:
        base = None
    if not base or len(base) < 6:
        return default
    try:
        r, g, b = int(base[0:2], 16), int(base[2:4], 16), int(base[4:6], 16)
    except Exception:
        return default
    transforms = []
    for child in el:
        t = etree.QName(child).localname
        v = child.get("val")
        if t in ("lumMod", "lumOff", "tint", "shade", "satMod") and v is not None:
            try:
                transforms.append((t, int(v) / 100000.0))
            except Exception:
                pass
    r, g, b = _apply_transforms(r, g, b, transforms)
    return "#%02X%02X%02X" % (min(255, max(0, r)), min(255, max(0, g)), min(255, max(0, b)))


def _fill_hex(fill, theme, default=None):
    """Hex of a python-pptx FillFormat's solid color, transforms included."""
    from pptx.oxml.ns import qn

    try:
        sf = fill._xPr.find(qn("a:solidFill"))
        if sf is not None and len(sf):
            return _clr_hex(sf[0], theme, default)
    except Exception:
        pass
    return default


def _run_hex(run, theme, default=None):
    """Hex of a run's explicit font color (rPr/solidFill), transforms included."""
    from pptx.oxml.ns import qn

    try:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is not None:
            sf = rPr.find(qn("a:solidFill"))
            if sf is not None and len(sf):
                return _clr_hex(sf[0], theme, default)
    except Exception:
        pass
    return default


def _is_dark(hexc):
    """Perceived-luminance test, to pick readable default text on a fill."""
    try:
        r, g, b = int(hexc[1:3], 16), int(hexc[3:5], 16), int(hexc[5:7], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 140
    except Exception:
        return False


def _slide_to_svg(slide, w_emu, h_emu, theme=None):
    """Render a slide's shapes to an SVG (scale-to-zero visual preview — no
    LibreOffice). EMU→px at 96dpi (EMU/9525); pt→px ×(96/72). Resolves theme
    colors, fills the true slide background, wraps text, and draws tables so the
    preview tracks the built deck instead of a cream-on-blank approximation."""
    import html as _html
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    theme = theme or {}
    U = 9525.0
    VW, VH = w_emu / U, h_emu / U

    # true slide background: walk slide → layout → master for a solid fill
    bg = "#FAF9F5"
    for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            if src.background.fill.type == 1:  # solid
                c = _fill_hex(src.background.fill, theme)
                if c:
                    bg = c
                    break
        except Exception:
            pass

    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW:.0f} {VH:.0f}" preserveAspectRatio="xMidYMid meet">']
    parts.append(f'<rect x="0" y="0" width="{VW:.0f}" height="{VH:.0f}" fill="{bg}"/>')

    for sh in slide.shapes:
        try:
            x, y = sh.left / U, sh.top / U
            w, h = sh.width / U, sh.height / U
        except Exception:
            continue

        # tables: draw the cell grid + text (data/comparison slides carry these)
        if getattr(sh, "has_table", False):
            try:
                tbl = sh.table
                col_ws = [c.width / U for c in tbl.columns] or [w]
                row_hs = [r.height / U for r in tbl.rows] or [h]
                sx, sy = w / (sum(col_ws) or w), h / (sum(row_hs) or h)
                col_ws = [c * sx for c in col_ws]
                row_hs = [r * sy for r in row_hs]
                ry = y
                for ri, rh in enumerate(row_hs):
                    rx = x
                    for ci, cw in enumerate(col_ws):
                        cell = tbl.cell(ri, ci)
                        cfill = _fill_hex(cell.fill, theme) or "#FFFFFF"
                        ctxt = "#F5F2EE" if _is_dark(cfill) else "#1A1917"
                        parts.append(f'<rect x="{rx:.1f}" y="{ry:.1f}" width="{cw:.1f}" height="{rh:.1f}" fill="{cfill}" stroke="#D9D5CC" stroke-width="0.5"/>')
                        ct = " ".join(cell.text.split())
                        if ct:
                            for li, wl in enumerate(_wrap(ct, max(1, int(cw / 6.5)))[:3]):
                                parts.append(f'<text x="{rx+4:.0f}" y="{ry+13+li*12:.0f}" font-family="Poppins,Arial,sans-serif" font-size="10" fill="{ctxt}">{_html.escape(wl)}</text>')
                        rx += cw
                    ry += rh
            except Exception:
                pass
            continue

        # fill (solid auto-shapes, incl. theme-colored full-bleed bands)
        fill = None
        try:
            if sh.fill.type is not None and sh.fill.type == 1:
                fill = _fill_hex(sh.fill, theme)
        except Exception:
            fill = None
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            parts.append(f'<rect x="{x:.0f}" y="{y:.0f}" width="{w:.0f}" height="{h:.0f}" fill="#E3DFD5"/>')
        elif sh.shape_type == MSO_SHAPE_TYPE.CHART:
            parts.append(f'<rect x="{x:.0f}" y="{y:.0f}" width="{w:.0f}" height="{h:.0f}" fill="#F1EEE7"/>')
            parts.append(f'<text x="{x+w/2:.0f}" y="{y+h/2:.0f}" font-family="Poppins,Arial" font-size="16" fill="#73716B" text-anchor="middle">chart</text>')
        elif fill:
            try:
                is_oval = sh.auto_shape_type == 9  # MSO_SHAPE.OVAL
            except Exception:
                is_oval = False
            if is_oval:
                parts.append(f'<ellipse cx="{x+w/2:.0f}" cy="{y+h/2:.0f}" rx="{w/2:.0f}" ry="{h/2:.0f}" fill="{fill}"/>')
            else:
                rx = 8 if (getattr(sh, "auto_shape_type", None) == 5) else 0
                parts.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" rx="{rx}" fill="{fill}"/>')

        # text, with word wrapping to the frame width
        if getattr(sh, "has_text_frame", False):
            tf = sh.text_frame
            anchor = tf.vertical_anchor
            lines = []  # (txt, px, bold, ital, col, fam, align, lh)
            for p in tf.paragraphs:
                runs = p.runs or []
                txt = "".join(r.text for r in runs)
                if not txt.strip():
                    continue
                sz = next((r.font.size.pt for r in runs if r.font.size), 18)
                bold = any(r.font.bold for r in runs)
                ital = any(r.font.italic for r in runs)
                # explicit run color, else pick readable default for the fill behind it
                col = next((_run_hex(r, theme) for r in runs if _run_hex(r, theme)), None)
                if not col:
                    col = "#F5F2EE" if (fill and _is_dark(fill)) else "#1A1917"
                fam = (next((r.font.name for r in runs if r.font.name), "Poppins") or "Poppins").split(",")[0]
                px = sz * 96 / 72
                lh = px * 1.22
                max_chars = max(1, int(w / (px * 0.52)))
                for wl in _wrap(txt, max_chars):
                    lines.append((wl, px, bold, ital, col, fam, p.alignment, lh))
            total_h = sum(ln[7] for ln in lines)
            if anchor == MSO_ANCHOR.MIDDLE:
                cy = y + (h - total_h) / 2
            elif anchor == MSO_ANCHOR.BOTTOM:
                cy = y + h - total_h
            else:
                cy = y
            for (txt, px, bold, ital, col, fam, align, lh) in lines:
                cy += lh
                if align == PP_ALIGN.CENTER:
                    tx, anc = x + w / 2, "middle"
                elif align == PP_ALIGN.RIGHT:
                    tx, anc = x + w, "end"
                else:
                    tx, anc = x, "start"
                weight = "700" if bold else "400"
                style = ' font-style="italic"' if ital else ""
                parts.append(
                    f'<text x="{tx:.0f}" y="{cy - lh*0.28:.0f}" font-family="{_html.escape(fam)},Arial,sans-serif" '
                    f'font-size="{px:.0f}" font-weight="{weight}" fill="{col}" text-anchor="{anc}"{style}>{_html.escape(txt)}</text>'
                )
    parts.append("</svg>")
    return "".join(parts)


def _slide_text(i, sl):
    """One slide as readable markdown — numbered so a reader can cite 'slide 4'."""
    head = f"## Slide {i + 1}" + (f": {sl['title']}" if sl["title"] else "")
    parts = [head]
    parts += [f"- {b}" for b in sl["bullets"]]
    for tbl in sl.get("tables", []):
        parts += [" | ".join(row) for row in tbl]
    for ch in sl.get("charts", []):
        parts.append(f"[{ch['type']} chart] categories: {', '.join(ch['categories'])}")
        for sr in ch["series"]:
            vals = ", ".join("—" if v is None else f"{v:g}" for v in sr["values"])
            parts.append(f"  {sr['name'] or 'series'}: {vals}")
    if sl.get("notes"):
        parts.append(f"Speaker notes: {sl['notes']}")
    return "\n".join(parts)


def _deck_design(prs):
    """Analyze the deck's LOOK & FEEL so the model can discuss design, not just
    text: dominant colors, fonts, slide size/aspect, and the layout mix. Pure
    python-pptx — no rendering."""
    from collections import Counter

    colors, fonts, sizes = Counter(), Counter(), Counter()
    pics = tables = charts = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.shape_type == 13:  # PICTURE
                pics += 1
            if getattr(sh, "has_table", False):
                tables += 1
            if getattr(sh, "has_chart", False):
                charts += 1
            try:
                if sh.fill.type == 1:
                    colors[f"#{sh.fill.fore_color.rgb}"] += 1
            except Exception:
                pass
            if getattr(sh, "has_text_frame", False):
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        # skip theme-token font refs (+mj-lt, +mn-lt) — not real names
                        if r.font.name and not r.font.name.startswith("+"):
                            fonts[r.font.name] += 1
                        try:
                            if r.font.color and r.font.color.type is not None:
                                colors[f"#{r.font.color.rgb}"] += 1
                        except Exception:
                            pass
                        if r.font.size:
                            sizes[round(r.font.size.pt)] += 1
    emu = 914400.0
    w, h = prs.slide_width / emu, prs.slide_height / emu
    aspect = "16:9" if abs(w / h - 16 / 9) < 0.05 else ("4:3" if abs(w / h - 4 / 3) < 0.05 else f"{w:.1f}x{h:.1f}in")
    return {
        "slide_count": len(prs.slides._sldIdLst),
        "aspect": aspect,
        "size_in": [round(w, 1), round(h, 1)],
        "palette": [c for c, _ in colors.most_common(6)],
        "fonts": [f for f, _ in fonts.most_common(4)],
        "font_sizes_pt": sorted({s for s, _ in sizes.most_common(6)}),
        "images": pics,
        "tables": tables,
        "charts": charts,
    }


def _design_text(d):
    lines = ["## Visual design (look & feel)"]
    lines.append(f"- Format: {d['aspect']} ({d['size_in'][0]}×{d['size_in'][1]} in), {d['slide_count']} slides")
    if d["palette"]:
        lines.append(f"- Color palette: {', '.join(d['palette'])}")
    if d["fonts"]:
        lines.append(f"- Fonts: {', '.join(d['fonts'])}")
    if d["font_sizes_pt"]:
        lines.append(f"- Type sizes: {', '.join(str(s) + 'pt' for s in d['font_sizes_pt'])}")
    visuals = []
    if d["images"]:
        visuals.append(f"{d['images']} images")
    if d["charts"]:
        visuals.append(f"{d['charts']} charts")
    if d["tables"]:
        visuals.append(f"{d['tables']} tables")
    if visuals:
        lines.append(f"- Visual elements: {', '.join(visuals)}")
    return "\n".join(lines)


def extract_preview(kind, file_bytes):
    """Structured preview extraction (pure-python, no LibreOffice). Returns
    {text, slides?, sheets?} so the client can render a readable preview in the
    scale-to-zero cloud where soffice/markitdown aren't available."""
    tmp = Path(tempfile.mkdtemp(prefix="preview-"))
    fp = tmp / f"in.{kind}"
    fp.write_bytes(file_bytes)
    out = {"kind": kind}
    if kind == "pptx":
        from pptx import Presentation

        prs = Presentation(str(fp))
        theme = _theme_colors(prs)
        slides, svgs = [], []
        for s in prs.slides:
            title, bullets = None, []
            for shape in s.shapes:
                if not shape.has_text_frame:
                    continue
                for i, para in enumerate(shape.text_frame.paragraphs):
                    txt = "".join(r.text for r in para.runs).strip()
                    if not txt:
                        continue
                    if title is None and shape == s.shapes.title:
                        title = txt
                    else:
                        bullets.append(txt)
            # tables and charts carry the substance on data slides — a deck read
            # without them comes back empty exactly where it matters most
            tables, charts = [], []
            for shape in s.shapes:
                if getattr(shape, "has_table", False):
                    tables.append([[c.text.strip() for c in row.cells] for row in shape.table.rows][:20])
                if getattr(shape, "has_chart", False):
                    try:
                        ch = shape.chart
                        cats = [str(c) for c in ch.plots[0].categories]
                        charts.append({
                            "type": str(ch.chart_type).split(" ")[0].lower(),
                            "categories": cats,
                            "series": [
                                {"name": sr.name or "", "values": [None if v is None else float(v) for v in sr.values]}
                                for sr in ch.series
                            ],
                        })
                    except Exception:
                        pass
            notes = ""
            try:
                if s.has_notes_slide:
                    notes = s.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
            slide = {"title": title or "", "bullets": bullets[:12]}
            if tables:
                slide["tables"] = tables
            if charts:
                slide["charts"] = charts
            if notes:
                slide["notes"] = notes
            slides.append(slide)
            try:
                svgs.append(_slide_to_svg(s, prs.slide_width, prs.slide_height, theme))
            except Exception:
                svgs.append(None)
        out["slides"] = slides
        out["svgs"] = svgs  # visual preview (rendered from shapes, no LibreOffice)
        design = _deck_design(prs)
        out["design"] = design
        text = "\n\n".join(_slide_text(i, sl) for i, sl in enumerate(slides))
        if design:
            text = f"{_design_text(design)}\n\n{text}"
        out["text"] = text
    elif kind == "docx":
        from docx import Document

        doc = Document(str(fp))
        blocks = []
        for p in doc.paragraphs:
            if p.text.strip():
                blocks.append({"style": p.style.name if p.style else "Normal", "text": p.text.strip()})
        for t in doc.tables:
            rows = [[c.text.strip() for c in r.cells] for r in t.rows]
            blocks.append({"style": "Table", "rows": rows[:20]})
        out["blocks"] = blocks[:200]
        # tables carry their rows into the text — a "[table]" placeholder reads
        # to the model as an empty table and it says so to the user
        out["text"] = "\n".join(
            b["text"] if "text" in b else "\n".join(" | ".join(r) for r in b.get("rows", []))
            for b in blocks
        )
    elif kind == "xlsx":
        from openpyxl import load_workbook

        # two passes: data_only gives cached values (present when Excel saved
        # the file), the raw pass gives formula text. A formula cell renders as
        # "=SUM(B2:B3) → 36" with the value, or just the formula without it —
        # data_only alone made formulas literally invisible (empty cells).
        wb = load_workbook(str(fp), data_only=True, read_only=True)
        wb_raw = load_workbook(str(fp), read_only=True)
        sheets = []
        for ws, ws_raw in zip(wb.worksheets, wb_raw.worksheets):
            rows = []
            for r, r_raw in zip(ws.iter_rows(values_only=True), ws_raw.iter_rows(values_only=True)):
                cells = []
                for val, raw in zip(r, r_raw):
                    if isinstance(raw, str) and raw.startswith("="):
                        cells.append(f"{raw} → {val}" if val is not None else raw)
                    else:
                        cells.append("" if val is None else str(val))
                rows.append(cells)
                if len(rows) >= 50:
                    break
            sheets.append({"name": ws.title, "rows": rows})
        out["sheets"] = sheets
        out["text"] = "\n\n".join(
            f"[{sh['name']}]\n" + "\n".join(" | ".join(r) for r in sh["rows"][:20]) for sh in sheets
        )
    elif kind == "pdf":
        import pdfplumber

        with pdfplumber.open(str(fp)) as pdf:
            out["text"] = "\n\n".join((pg.extract_text() or "") for pg in pdf.pages[:30])
    else:
        out["text"] = ""
    return {"ok": True, **out}


def _source_bytes(event):
    """Extract source: inline base64, or an S3 object. A synchronous invoke caps
    the request at 6MB and base64 inflates by 4/3, so callers hand anything
    larger over as a key — this function shares the app's role and can read the
    uploads bucket itself."""
    if event.get("file_b64"):
        return base64.b64decode(event["file_b64"])
    src = event.get("s3")
    if src and src.get("bucket") and src.get("key"):
        import boto3

        return boto3.client("s3").get_object(Bucket=src["bucket"], Key=src["key"])["Body"].read()
    return None


def handler(event, _context):
    # preview extraction path (no build)
    if event.get("op") == "extract":
        kind = event.get("kind")
        if kind not in {"pptx", "docx", "xlsx", "pdf"}:
            return {"ok": False, "error": f"cannot preview kind: {kind}"}
        try:
            data = _source_bytes(event)
        except Exception as err:  # noqa: BLE001
            return {"ok": False, "error": f"source unreadable: {type(err).__name__}: {err}"}
        if not data:
            return {"ok": False, "error": "no file_b64 or s3 source given"}
        try:
            return extract_preview(kind, data)
        except Exception as err:  # noqa: BLE001
            return {"ok": False, "error": f"{type(err).__name__}: {err}"}

    skill = event.get("skill")
    payload = event.get("payload")
    if skill not in SKILLS or payload is None:
        return {"ok": False, "error": f"unsupported skill: {skill}"}

    tmp = Path(tempfile.mkdtemp(prefix="office-"))
    payload_file = tmp / "payload.json"
    payload_file.write_text(json.dumps(payload))
    ext = "pptx" if skill == "pptx" else skill
    out_file = tmp / f"out.{ext}"

    argv = ["build", "--payload", str(payload_file), "--out", str(out_file)]
    tmpl = HERE / "templates" / TEMPLATES.get(skill, "")
    if skill in TEMPLATES and tmpl.exists():
        argv += ["--template", str(tmpl)]

    mod = importlib.import_module(f"build_{skill}")
    old_argv, old_stdout, old_stderr = sys.argv, sys.stdout, sys.stderr
    sys.argv = argv
    captured = tempfile.SpooledTemporaryFile(mode="w+")
    captured_err = tempfile.SpooledTemporaryFile(mode="w+")
    sys.stdout = captured
    sys.stderr = captured_err
    result = {"ok": False, "error": "builder produced no output"}
    try:
        try:
            mod.main()
        except SystemExit:
            pass
        captured.seek(0)
        last = [ln for ln in captured.read().strip().splitlines() if ln.strip()]
        if last:
            result = json.loads(last[-1])
        else:
            # The builder exited via fail() without a result line — its specific
            # reason (e.g. "spec validation failed (1): slide 2: 50 words on a
            # content slide") went to stderr. Surface THAT, not a generic
            # "produced no output", so the app's fix-and-rerender loop gets
            # actionable findings to regenerate against.
            captured_err.seek(0)
            errlines = [ln for ln in captured_err.read().strip().splitlines() if ln.strip()]
            if errlines:
                result = {"ok": False, "error": errlines[-1]}
    except Exception as err:  # noqa: BLE001
        result = {"ok": False, "error": f"{type(err).__name__}: {err}"}
    finally:
        sys.argv, sys.stdout, sys.stderr = old_argv, old_stdout, old_stderr

    if out_file.exists():
        result["file_b64"] = base64.b64encode(out_file.read_bytes()).decode()
    return result


# local smoke: python lambda_handler.py '{"skill":"xlsx","payload":{...}}'
if __name__ == "__main__":
    print(json.dumps(handler(json.loads(sys.argv[1]), None))[:200])

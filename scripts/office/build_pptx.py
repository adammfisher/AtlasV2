#!/usr/bin/env python3
"""Deterministic PPTX designer (skills/pptx/schema.json → the corporate template).

Renders the 12 layout archetypes onto dfs_default.potx: title text lands in the
template's named title placeholders (normalized to the 12-column grid), content
frames are grid-computed per archetype — the MODEL never positions anything
(frontier-tier position_overrides excepted). All color styles through MSO theme
slots via pptx_design roles; text is measured per line with embedded Poppins
metrics and stepped down within its doctrine range; a frame still overflowing
at the 14pt floor raises an OVERFLOW flag for the validator — never silently
accepted.
"""
import os
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

import pptx_design as D
import validate_common as vc

# doctrine type scale: (max_pt, min_pt) per role
SCALE = {
    "deck_title": (40, 32),
    "title": (36, 28),
    "section": (24, 20),
    "body": (24, 18),
    "caption": (14, 12),
    "stat": (72, 60),
    "quote": (32, 28),
}
FOOTER_PT = 12
CONTENT_TOP_GAP = 0.5   # whitespace below the headline (>= 0.5")
FOOTER_Y = 6.85
CONTENT_BOTTOM = 6.45   # >= 0.3" clearance above the footer line
STATEMENT = {"title", "section_divider", "quote", "closing_cta"}

# Furniture-free DFS layouts only (stock layouts 0-8 carry decorative banner
# rectangles that collide with grid content; 12-19 are clean, most with a TITLE
# placeholder; the master supplies the brand logo bottom-left on every slide).
LAYOUT_IDX = {
    "title": 16,             # Title Slide - Sub heading (clean, TITLE ph)
    "agenda": 17,            # 1_Title Only (clean, TITLE ph)
    "section_divider": 14,   # 2_Custom Layout (clean, TITLE ph)
    "content_bullets": 17,
    "content_chart": 13,     # chart-left 2-col layout (clean, TITLE ph)
    "comparison": 14,
    "big_stat": 17,
    "quote": 6,              # Blank (statement field; no empty ph left behind)
    "timeline_process": 14,
    "two_column": 15,        # 1_Custom Layout (clean, TITLE ph)
    "table": 14,
    "closing_cta": 20,       # Close (brand logo lockup center)
}


class Ctx:
    def __init__(self, prs, deck_title, total):
        self.prs = prs
        self.deck_title = deck_title
        self.total = total
        self.overflows = []      # [{slide, frame, detail}] — validator input
        self.keyed = {}          # (slide_idx, key) -> shape, for overrides


# ── low-level styled primitives (theme slots only) ─────────────────────────
def _font(run_font, role, size_pt, bold=False, italic=False, brightness=0.0):
    run_font.size = Pt(size_pt)
    run_font.bold = bold
    run_font.italic = italic
    run_font.name = D.BODY
    run_font.color.theme_color = D.ROLES[role]
    if brightness:
        run_font.color.brightness = brightness


def _fill(shape, role, brightness=0.0):
    shape.fill.solid()
    shape.fill.fore_color.theme_color = D.ROLES[role]
    if brightness:
        shape.fill.fore_color.brightness = brightness
    shape.line.fill.background()
    shape.shadow.inherit = False


def _bg(slide, role):
    slide.background.fill.solid()
    slide.background.fill.fore_color.theme_color = D.ROLES[role]


def _frame(shape, x, y, w, h):
    shape.left, shape.top = Emu(D.emu(x)), Emu(D.emu(y))
    shape.width, shape.height = Emu(D.emu(w)), Emu(D.emu(h))


def _textbox(ctx, slide, sidx, key, x, y, w, h, text, role, scale_key, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             brightness=0.0, spacing=1.15, max_lines=None):
    """Measured text: fit within the doctrine range, flag overflow, never guess."""
    max_pt, min_pt = SCALE[scale_key]
    fit = D.fit_text(text, w, h, max_pt, min_pt, bold=bold, spacing=spacing)
    if fit["overflow"]:
        ctx.overflows.append({"slide": sidx, "frame": key, "detail": f"overflow at {fit['size']}pt floor"})
    if max_lines and fit["lines"] > max_lines:
        ctx.overflows.append({"slide": sidx, "frame": key, "detail": f"{fit['lines']} lines (max {max_lines})"})
    tb = slide.shapes.add_textbox(Emu(D.emu(x)), Emu(D.emu(y)), Emu(D.emu(w)), Emu(D.emu(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = spacing
    run = para.add_run()
    run.text = str(text)
    _font(run.font, role, fit["size"], bold=bold, italic=italic, brightness=brightness)
    ctx.keyed[(sidx, key)] = tb
    return tb, fit


def _title_into_placeholder(ctx, slide, sidx, spec, dark=False, y=D.MARGIN,
                            scale_key="title", w_cols=12, align=PP_ALIGN.LEFT):
    """The layout's TITLE placeholder, normalized to the grid. Left-aligned,
    stepped within its scale range, <= 2 measured lines. NO accent line
    beneath — ever. Falls back to a textbox on placeholder-less layouts."""
    ph = next((p for p in slide.placeholders if p.placeholder_format.idx == 0), None)
    w = D.grid_w(w_cols)
    x = D.MARGIN if align == PP_ALIGN.LEFT else (D.SLIDE_W - w) / 2
    fit = D.fit_text(spec["title"], w, 2.2, *SCALE[scale_key], bold=True, spacing=1.05)
    if fit["lines"] > 2:
        ctx.overflows.append({"slide": sidx, "frame": "title", "detail": f"title wraps to {fit['lines']} lines"})
    h = max(0.6, fit["lines"] * D.line_height_in(fit["size"], bold=True, spacing=1.05))
    role = "text_inverse" if dark else "text"
    if ph is None:
        _textbox(ctx, slide, sidx, "title", x, y, w, h + 0.1, spec["title"],
                 role, scale_key, bold=True, spacing=1.05, max_lines=2, align=align)
        return y + h
    _frame(ph, x, y, w, h + 0.1)
    tf = ph.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = 1.05
    run = para.runs[0] if para.runs else para.add_run()
    run.text = str(spec["title"])
    _font(run.font, role, fit["size"], bold=True)
    ctx.keyed[(sidx, "title")] = ph
    return y + h


def _footer(ctx, slide, sidx):
    # bottom-left belongs to the master's brand logo — page number only, right
    right = slide.shapes.add_textbox(
        Emu(D.emu(D.SLIDE_W - D.MARGIN - 1.5)), Emu(D.emu(FOOTER_Y)), Emu(D.emu(1.5)), Emu(D.emu(0.3)))
    para = right.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = f"{sidx} / {ctx.total}"
    _font(run.font, "text", FOOTER_PT, brightness=0.35)


def _notes(slide, spec):
    slide.notes_slide.notes_text_frame.text = str(spec.get("speaker_notes", ""))


def _bullets_block(ctx, slide, sidx, items, x, w, top, bottom, icons=None):
    """Bullet stack sized collectively: the largest body size where every item
    fits its row. Icon column (theme-shape glyphs) when icons are given."""
    n = max(1, len(items))
    avail = bottom - top
    row_h = min(1.1, avail / n)
    # the text column is narrower than the block: icon column or bullet marker
    text_x, text_w = (x + 0.75, w - 0.75) if icons else (x + 0.4, w - 0.4)
    max_pt, min_pt = SCALE["body"]
    size = max_pt
    while size > 14:
        if all(D.required_height_in(b, text_w, size) <= row_h - 0.12 for b in items):
            break
        size -= 1
    if size < min_pt:
        # below the projected-body range: legal only above the 14pt floor; flag either way
        ctx.overflows.append({"slide": sidx, "frame": "bullets",
                              "detail": f"body stepped to {size}pt (range floor {min_pt}pt)"})
    for i, item in enumerate(items):
        y = top + i * row_h
        if icons:
            _icon(slide, icons[i % len(icons)], x, y + 0.02, 0.5)
        else:
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(D.emu(x + 0.02)), Emu(D.emu(y + 0.11)), Emu(D.emu(0.12)), Emu(D.emu(0.12)))
            _fill(marker, "accent")
        tb = slide.shapes.add_textbox(
            Emu(D.emu(text_x)), Emu(D.emu(y)), Emu(D.emu(text_w)), Emu(D.emu(row_h)))
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.line_spacing = 1.1
        run = para.add_run()
        run.text = str(item)
        _font(run.font, "text", size)
    ctx.keyed[(sidx, "body")] = tb
    return size


# icon primitives: name -> [(shape, dx, dy, dw, dh, role, brightness)]
_ICONS = {
    "growth": [(MSO_SHAPE.RECTANGLE, 0.02, 0.30, 0.10, 0.18, "accent", 0.0),
               (MSO_SHAPE.RECTANGLE, 0.18, 0.18, 0.10, 0.30, "accent", 0.0),
               (MSO_SHAPE.RECTANGLE, 0.34, 0.04, 0.10, 0.44, "supporting", 0.0)],
    "target": [(MSO_SHAPE.OVAL, 0.02, 0.02, 0.44, 0.44, "accent", 0.0),
               (MSO_SHAPE.OVAL, 0.15, 0.15, 0.18, 0.18, "background", 0.0)],
    "people": [(MSO_SHAPE.OVAL, 0.06, 0.02, 0.16, 0.16, "supporting", 0.0),
               (MSO_SHAPE.OVAL, 0.26, 0.02, 0.16, 0.16, "accent", 0.0),
               (MSO_SHAPE.ROUNDED_RECTANGLE, 0.04, 0.22, 0.40, 0.24, "supporting", 0.6)],
    "risk": [(MSO_SHAPE.ISOSCELES_TRIANGLE, 0.02, 0.02, 0.44, 0.42, "accent", 0.0),
             (MSO_SHAPE.RECTANGLE, 0.21, 0.14, 0.06, 0.16, "background", 0.0)],
    "time": [(MSO_SHAPE.OVAL, 0.02, 0.02, 0.44, 0.44, "supporting", 0.0),
             (MSO_SHAPE.RECTANGLE, 0.22, 0.10, 0.05, 0.16, "background", 0.0),
             (MSO_SHAPE.RECTANGLE, 0.22, 0.22, 0.14, 0.05, "background", 0.0)],
    "money": [(MSO_SHAPE.OVAL, 0.02, 0.02, 0.44, 0.44, "accent", 0.0),
              (MSO_SHAPE.RECTANGLE, 0.21, 0.10, 0.06, 0.28, "background", 0.0)],
    "check": [(MSO_SHAPE.ROUNDED_RECTANGLE, 0.02, 0.02, 0.44, 0.44, "accent", 0.0),
              (MSO_SHAPE.CHEVRON, 0.12, 0.16, 0.24, 0.16, "background", 0.0)],
    "idea": [(MSO_SHAPE.OVAL, 0.08, 0.02, 0.32, 0.32, "accent", 0.0),
             (MSO_SHAPE.RECTANGLE, 0.18, 0.36, 0.12, 0.10, "supporting", 0.0)],
    "globe": [(MSO_SHAPE.OVAL, 0.02, 0.02, 0.44, 0.44, "supporting", 0.0),
              (MSO_SHAPE.RECTANGLE, 0.02, 0.22, 0.44, 0.04, "background", 0.0)],
    "gear": [(MSO_SHAPE.OVAL, 0.08, 0.08, 0.32, 0.32, "supporting", 0.0),
             (MSO_SHAPE.OVAL, 0.18, 0.18, 0.12, 0.12, "background", 0.0)],
}


def _icon(slide, name, x, y, box):
    for shape_kind, dx, dy, dw, dh, role, bright in _ICONS.get(name, _ICONS["target"]):
        scale = box / 0.48
        sh = slide.shapes.add_shape(
            shape_kind, Emu(D.emu(x + dx * scale)), Emu(D.emu(y + dy * scale)),
            Emu(D.emu(dw * scale)), Emu(D.emu(dh * scale)))
        _fill(sh, role, bright)


def _sorted_chart(cs):
    cats = [str(c) for c in cs["categories"]]
    series = [{"name": s["name"], "values": list(s["values"])} for s in cs["series"]]
    if cs.get("sort") in ("value_desc", "value_asc") and series:
        order = sorted(range(len(cats)), key=lambda i: series[0]["values"][i],
                       reverse=cs["sort"] == "value_desc")
        cats = [cats[i] for i in order]
        for s in series:
            s["values"] = [s["values"][i] for i in order]
    return cats, series


def _chart(ctx, slide, sidx, cs, x, y, w, h):
    cats, series = _sorted_chart(cs)
    data = CategoryChartData()
    data.categories = cats
    for s in series[:5]:
        data.add_series(str(s["name"]), s["values"])
    kind = {"line": XL_CHART_TYPE.LINE_MARKERS, "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE}[cs["kind"]]
    gf = slide.shapes.add_chart(kind, Emu(D.emu(x)), Emu(D.emu(y)), Emu(D.emu(w)), Emu(D.emu(h)), data)
    chart = gf.chart
    chart.has_title = False
    for i, plot_series in enumerate(chart.series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.theme_color = D.CHART_SERIES_SLOTS[i % len(D.CHART_SERIES_SLOTS)]
        plot_series.format.line.fill.background()
    if cs["kind"] != "pie":
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.color.theme_color = D.ROLES["text"]
        try:
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.theme_color = D.ROLES["supporting"]
            chart.value_axis.major_gridlines.format.line.color.brightness = 0.85
            chart.category_axis.has_major_gridlines = False
            for axis in (chart.value_axis, chart.category_axis):
                axis.tick_labels.font.size = Pt(12)
                axis.tick_labels.font.color.theme_color = D.ROLES["text"]
        except Exception:
            pass  # pie/axis-less chart types
    multi = len(series) > 1
    chart.has_legend = multi or cs["kind"] == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.theme_color = D.ROLES["text"]
    ctx.keyed[(sidx, "chart")] = gf
    return gf


# ── archetype renderers ────────────────────────────────────────────────────
def r_title(ctx, slide, sidx, spec):
    _bg(slide, "dominant_dark")
    # hero on the lower-third power line (rule of thirds)
    _title_into_placeholder(ctx, slide, sidx, spec, dark=True, y=2.6,
                            scale_key="deck_title", w_cols=9)
    if spec.get("subtitle"):
        _textbox(ctx, slide, sidx, "subtitle", D.MARGIN, 5.0, D.grid_w(8), 0.8,
                 spec["subtitle"], "text_inverse", "body", brightness=-0.15)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(D.emu(D.grid_x(10))), Emu(D.emu(5.0)), Emu(D.emu(1.4)), Emu(D.emu(1.4)))
    _fill(accent, "accent")


def r_agenda(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    items = spec["bullets"][:6]
    row_h = min(0.8, (CONTENT_BOTTOM - top) / max(1, len(items)))
    for i, item in enumerate(items):
        y = top + i * row_h
        num = slide.shapes.add_textbox(Emu(D.emu(D.grid_x(0))), Emu(D.emu(y)), Emu(D.emu(0.6)), Emu(D.emu(row_h)))
        para = num.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{i + 1:02d}"
        _font(run.font, "supporting", 20, bold=True)
        tb = slide.shapes.add_textbox(
            Emu(D.emu(D.grid_x(1))), Emu(D.emu(y)), Emu(D.emu(D.grid_w(7))), Emu(D.emu(row_h)))
        para = tb.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(item)
        _font(run.font, "text", 20)
    # cols 9–12 stay empty by design (whitespace budget)
    _footer(ctx, slide, sidx)


def r_section_divider(ctx, slide, sidx, spec):
    _bg(slide, "dominant_dark")
    num = slide.shapes.add_textbox(Emu(D.emu(D.MARGIN)), Emu(D.emu(0.8)), Emu(D.emu(2.5)), Emu(D.emu(1.0)))
    para = num.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = f"{sidx:02d}"
    _font(run.font, "text_inverse", 24, bold=True, brightness=-0.3)
    _title_into_placeholder(ctx, slide, sidx, spec, dark=True, y=4.4,
                            scale_key="deck_title", w_cols=10)
    if spec.get("subtitle"):
        _textbox(ctx, slide, sidx, "subtitle", D.MARGIN, 6.1, D.grid_w(8), 0.7,
                 spec["subtitle"], "text_inverse", "body", brightness=-0.15)


def r_content_bullets(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    icons = spec.get("icons")
    _bullets_block(ctx, slide, sidx, spec["bullets"], D.grid_x(0), D.grid_w(8), top, CONTENT_BOTTOM,
                   icons=icons)
    if not icons:  # visual element: a quiet accent mark keeps the slide from being text-only
        band = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(D.emu(D.grid_x(9))), Emu(D.emu(top + 0.1)),
            Emu(D.emu(D.grid_w(3))), Emu(D.emu(CONTENT_BOTTOM - top - 0.2)))
        _fill(band, "supporting", D.PANEL_BRIGHTNESS)
        _icon(slide, "target", D.grid_x(10) + 0.2, (top + CONTENT_BOTTOM) / 2 - 0.4, 0.8)
    _footer(ctx, slide, sidx)


def r_content_chart(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    bullets = spec.get("bullets") or []
    chart_w = D.grid_w(8 if bullets else 12)
    _chart(ctx, slide, sidx, spec["chart"], D.grid_x(0), top, chart_w, CONTENT_BOTTOM - top)
    if bullets:
        _bullets_block(ctx, slide, sidx, bullets[:3], D.grid_x(9), D.grid_w(3), top + 0.2, CONTENT_BOTTOM)
    _footer(ctx, slide, sidx)


def r_comparison(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    cols = spec["columns"]
    n = len(cols)
    span = 12 // n
    for c, col in enumerate(cols):
        x = D.grid_x(c * span)
        w = D.grid_w(span) - (D.BLOCK_GAP if c < n - 1 else 0)
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(D.emu(x)), Emu(D.emu(top)),
            Emu(D.emu(w)), Emu(D.emu(CONTENT_BOTTOM - top)))
        _fill(panel, "supporting", D.PANEL_BRIGHTNESS)
        head = slide.shapes.add_textbox(
            Emu(D.emu(x + 0.3)), Emu(D.emu(top + 0.25)), Emu(D.emu(w - 0.6)), Emu(D.emu(0.5)))
        para = head.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(col["head"])
        _font(run.font, "supporting", 20, bold=True)
        yy = top + 0.95
        for item in col["items"]:
            fit = D.fit_text(item, w - 0.9, 0.9, 18, 14, spacing=1.08)
            if fit["overflow"]:
                ctx.overflows.append({"slide": sidx, "frame": f"column {c + 1}", "detail": "item overflow"})
            # frame height = measured height (fixed heights taller than the
            # stacking step collide — the visual gate caught exactly that)
            item_h = D.required_height_in(item, w - 0.9, fit["size"], spacing=1.08) + 0.06
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(D.emu(x + 0.3)), Emu(D.emu(yy + 0.1)), Emu(D.emu(0.1)), Emu(D.emu(0.1)))
            _fill(marker, "accent")
            tb = slide.shapes.add_textbox(
                Emu(D.emu(x + 0.55)), Emu(D.emu(yy)),
                Emu(D.emu(w - 0.9)), Emu(D.emu(item_h)))
            tf = tb.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.line_spacing = 1.08
            run = para.add_run()
            run.text = str(item)
            _font(run.font, "text", fit["size"])
            yy += item_h + 0.12
    _footer(ctx, slide, sidx)


def r_big_stat(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    stat = spec["stat"]
    # numeral on the left-third power point; >= 40% of the slide stays empty
    _textbox(ctx, slide, sidx, "stat", D.grid_x(0), top + 0.4, D.grid_w(7), 1.6,
             stat["value"], "text", "stat", bold=True, spacing=1.0)
    mark = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(D.emu(D.grid_x(0))), Emu(D.emu(top + 2.15)), Emu(D.emu(1.1)), Emu(D.emu(0.12)))
    _fill(mark, "accent")
    _textbox(ctx, slide, sidx, "stat_label", D.grid_x(0), top + 2.45, D.grid_w(7), 0.6,
             stat["label"], "text", "section", bold=False, brightness=0.25)
    if spec.get("support"):
        _textbox(ctx, slide, sidx, "support", D.grid_x(0), top + 3.15, D.grid_w(8), 0.8,
                 spec["support"], "text", "body")
    _footer(ctx, slide, sidx)


def r_quote(ctx, slide, sidx, spec):
    _bg(slide, "dominant_dark")
    # glyph frame sized to its measured 96pt line height, clear of the quote text
    glyph_h = D.line_height_in(96, bold=True, spacing=1.0) + 0.1
    glyph = slide.shapes.add_textbox(Emu(D.emu(D.MARGIN)), Emu(D.emu(0.55)), Emu(D.emu(1.5)), Emu(D.emu(glyph_h)))
    para = glyph.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "“"
    _font(run.font, "accent", 96, bold=True)
    _textbox(ctx, slide, sidx, "quote", D.grid_x(1), 2.5, D.grid_w(10), 2.4, spec["quote"],
             "text_inverse", "quote", italic=True, spacing=1.2)
    _textbox(ctx, slide, sidx, "attribution", D.grid_x(1), 5.4, D.grid_w(10), 0.5,
             f"— {spec['attribution']}", "text_inverse", "caption", brightness=-0.15,
             align=PP_ALIGN.RIGHT)


def r_timeline(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    steps = spec["steps"]
    n = len(steps)
    mid = top + (CONTENT_BOTTOM - top) / 2 - 0.6
    track = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(D.emu(D.grid_x(0) + 0.3)), Emu(D.emu(mid + 0.28)),
        Emu(D.emu(D.grid_w(12) - 0.6)), Emu(D.emu(0.03)))
    _fill(track, "supporting", 0.6)
    span = 12 // n if n else 12
    for i, step in enumerate(steps):
        x = D.grid_x(i * span)
        w = D.grid_w(span) - 0.2
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(D.emu(x + 0.15)), Emu(D.emu(mid)), Emu(D.emu(0.6)), Emu(D.emu(0.6)))
        _fill(node, "accent")
        num = node.text_frame
        para = num.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = str(i + 1)
        _font(run.font, "text_inverse", 18, bold=True)  # >=14pt bold on teal: 3:1 large-text bar
        label_fit = D.fit_text(step["label"], w, 0.9, 18, 14, bold=True, spacing=1.05)
        label_h = D.required_height_in(step["label"], w, label_fit["size"], bold=True, spacing=1.05) + 0.06
        label = slide.shapes.add_textbox(
            Emu(D.emu(x)), Emu(D.emu(mid + 0.85)), Emu(D.emu(w)), Emu(D.emu(label_h)))
        label.text_frame.word_wrap = True
        para = label.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(step["label"])
        _font(run.font, "text", label_fit["size"], bold=True)
        if step.get("detail"):
            detail = slide.shapes.add_textbox(
                Emu(D.emu(x)), Emu(D.emu(mid + 0.95 + label_h)), Emu(D.emu(w)), Emu(D.emu(0.85)))
            tf = detail.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.line_spacing = 1.08
            run = para.add_run()
            run.text = str(step["detail"])
            _font(run.font, "text", 14, brightness=0.2)
    _footer(ctx, slide, sidx)


def r_two_column(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    _bullets_block(ctx, slide, sidx, spec["bullets"], D.grid_x(0), D.grid_w(6) - 0.2, top, CONTENT_BOTTOM)
    if spec.get("chart"):
        _chart(ctx, slide, sidx, spec["chart"], D.grid_x(6), top, D.grid_w(6), CONTENT_BOTTOM - top)
    elif spec.get("stat"):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(D.emu(D.grid_x(6))), Emu(D.emu(top)),
            Emu(D.emu(D.grid_w(6))), Emu(D.emu(CONTENT_BOTTOM - top)))
        _fill(panel, "supporting", D.PANEL_BRIGHTNESS)
        _textbox(ctx, slide, sidx, "stat", D.grid_x(6) + 0.4, top + 0.8, D.grid_w(6) - 0.8, 1.4,
                 spec["stat"]["value"], "text", "stat", bold=True)
        _textbox(ctx, slide, sidx, "stat_label", D.grid_x(6) + 0.4, top + 2.4, D.grid_w(6) - 0.8, 0.8,
                 spec["stat"]["label"], "text", "body", brightness=0.2)
    else:
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(D.emu(D.grid_x(6))), Emu(D.emu(top)),
            Emu(D.emu(D.grid_w(6))), Emu(D.emu(CONTENT_BOTTOM - top)))
        _fill(panel, "supporting", D.PANEL_BRIGHTNESS)
        names = spec.get("icons") or ["target", "growth", "check"]
        for i, name in enumerate(names[:3]):
            _icon(slide, name, D.grid_x(7) + i * 1.3, top + (CONTENT_BOTTOM - top) / 2 - 0.4, 0.8)
    _footer(ctx, slide, sidx)


def r_table(ctx, slide, sidx, spec):
    _bg(slide, "background")
    top = _title_into_placeholder(ctx, slide, sidx, spec) + CONTENT_TOP_GAP
    headers = spec["table"]["headers"]
    rows = spec["table"]["rows"]
    ncol = len(headers)
    shape = slide.shapes.add_table(
        1 + len(rows), ncol, Emu(D.emu(D.grid_x(0))), Emu(D.emu(top)),
        Emu(D.emu(D.grid_w(12))), Emu(D.emu(min(CONTENT_BOTTOM - top, 0.5 * (1 + len(rows))))))
    table = shape.table
    table.first_row = True
    table.horz_banding = True
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.theme_color = D.ROLES["supporting"]
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                _font(run.font, "text_inverse", 14, bold=True)
    for r, row in enumerate(rows, start=1):
        for c in range(ncol):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.theme_color = D.ROLES["background"]
            if r % 2 == 0:
                cell.fill.fore_color.theme_color = D.ROLES["supporting"]
                cell.fill.fore_color.brightness = D.BAND_BRIGHTNESS
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    _font(run.font, "text", 14)
    ctx.keyed[(sidx, "table")] = shape
    _footer(ctx, slide, sidx)


def r_closing(ctx, slide, sidx, spec):
    # the Close layout carries the brand logo lockup center (2.89–4.64) —
    # closing message sits above it, CTA/contact below; no extra decoration
    _bg(slide, "dominant_dark")
    _textbox(ctx, slide, sidx, "title", D.grid_x(1), 1.1, D.grid_w(10), 1.5, spec["title"],
             "text_inverse", "deck_title", bold=True, align=PP_ALIGN.CENTER, max_lines=2)
    if spec.get("subtitle"):
        _textbox(ctx, slide, sidx, "subtitle", D.grid_x(2), 5.0, D.grid_w(8), 0.8,
                 spec["subtitle"], "text_inverse", "body", align=PP_ALIGN.CENTER, brightness=-0.15)


RENDERERS = {
    "title": r_title, "agenda": r_agenda, "section_divider": r_section_divider,
    "content_bullets": r_content_bullets, "content_chart": r_content_chart,
    "comparison": r_comparison, "big_stat": r_big_stat, "quote": r_quote,
    "timeline_process": r_timeline, "two_column": r_two_column, "table": r_table,
    "closing_cta": r_closing,
}


def _apply_overrides(ctx, sidx, spec):
    for ov in spec.get("position_overrides") or []:
        shape = ctx.keyed.get((sidx, ov["shape"]))
        if shape is not None:
            _frame(shape, ov["x_in"], ov["y_in"], ov["w_in"], ov["h_in"])


def _assert_theme_only():
    """No RGB literal outside the palette module: this builder and the design
    module must never import python-pptx's RGB color class."""
    # needles assembled in halves so this function's own source never matches
    import_needle = "pptx.dml" + ".color"
    class_needle = "RGB" + "Color"
    here = Path(__file__).resolve().parent
    for name in ("build_pptx.py", "pptx_design.py"):
        src = (here / name).read_text()
        if import_needle in src or class_needle in src:
            raise AssertionError(f"{name}: RGB color literal found — style through pptx_design roles only")


def build(payload: dict, template: str, out: Path) -> dict:
    _assert_theme_only()
    template_file = template or str(
        Path(__file__).resolve().parents[2] / "skills/pptx/templates/dfs_default.potx")
    if not Path(template_file).exists():  # bundled Lambda layout
        template_file = str(Path(__file__).resolve().parent / "templates/dfs_default.potx")
    prs = Presentation(template_file)
    prs.slide_width = Emu(D.emu(D.SLIDE_W))
    prs.slide_height = Emu(D.emu(D.SLIDE_H))
    # a .potx opens with zero slides; render every spec slide onto its layout
    slides = payload["slides"]
    ctx = Ctx(prs, payload.get("title", ""), len(slides))
    for i, spec in enumerate(slides, start=1):
        layout = prs.slide_layouts[LAYOUT_IDX.get(spec["archetype"], 5)]
        slide = prs.slides.add_slide(layout)
        # drop layout-inherited placeholders we don't fill (date/footer/number stubs)
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx != 0:
                ph._element.getparent().remove(ph._element)
        RENDERERS[spec["archetype"]](ctx, slide, i, spec)
        # an unfilled title placeholder would render as a "click to edit" stub
        for ph in list(slide.placeholders):
            if not ph.text_frame.text.strip():
                ph._element.getparent().remove(ph._element)
        _apply_overrides(ctx, i, spec)
        _notes(slide, spec)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {
        "slides": len(slides),
        "bytes": out.stat().st_size,
        "overflow_flags": ctx.overflows,
    }


def extract_texts(path: Path) -> list:
    texts = []
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return texts


def main() -> None:
    args = vc.cli("pptx builder")
    payload = vc.load_payload(args.payload)
    vc.spec_gate("pptx", payload)
    out = Path(args.out)
    meta = build(payload, args.template or "", out)

    checks = [vc.openxml_audit(out), vc.zip_sanity(out)]
    reopened = Presentation(str(out))
    texts = extract_texts(out)
    checks.append(
        vc.check("Round-trip", len(list(reopened.slides)) == meta["slides"] and any(t.strip() for t in texts))
    )
    checks.append(vc.check("Overflow-free (measured)", not meta["overflow_flags"]))
    # THE HARD GATE — overflow, collision, margins, contrast, fonts, content,
    # placeholders, speaker notes on the BUILT file. Findings feed the server's
    # bounded fix-and-rerender loop; a failing deck is never a success.
    findings = vc.visual_gate_pptx(out, args.template or None)
    checks.append(vc.check("Visual gate (deterministic)", not findings))
    checks.append(vc.post_render_bleed(out))
    meta["findings"] = ([f"OVERFLOW slide {o['slide']} {o['frame']}: {o['detail']}" for o in meta["overflow_flags"]] + findings)[:12]
    meta["overflow_flags"] = len(meta["overflow_flags"])
    if os.environ.get("ATLAS_VISION_CRITIQUE") == "1":
        meta["thumbs_b64"] = vc.render_thumbnails(out)
    vc.emit(out, meta, checks)


if __name__ == "__main__":
    try:
        main()
    except SystemExit:
        raise
    except Exception as err:  # helper contract: stderr + exit 1
        vc.fail(f"{type(err).__name__}: {err}")

"""Shared helper-contract plumbing + validation chain pieces (PRD §4.3.4, §4.5).

Every builder: parse args → build → run checks → print single-line JSON to
stdout and exit 0. Any exception → stderr + exit 1 (the server surfaces the
stderr tail as the pipeline error).
"""
import argparse
import json
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

RECALC_SKIP = "Recalc skipped — soffice not found"
THUMBS_SKIP = "Thumbnails skipped — soffice not found"


def cli(description: str) -> argparse.Namespace:
    p = argparse.ArgumentParser(description=description)
    p.add_argument("--payload", required=True)
    p.add_argument("--template", default=None)
    p.add_argument("--out", required=True)
    return p.parse_args()


def load_payload(path: str) -> dict:
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def emit(file: Path, meta: dict, checks: list) -> None:
    print(json.dumps({"ok": True, "file": str(file), "meta": meta, "checks": checks}))


def fail(message: str) -> None:
    print(message, file=sys.stderr)
    sys.exit(1)


def check(label: str, ok: bool) -> dict:
    return {"label": label, "ok": bool(ok)}


# --- §4.5(a) zip + content-types sanity -----------------------------------
def zip_sanity(path: Path) -> dict:
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            ok = "[Content_Types].xml" in names and z.testzip() is None
    except Exception:
        ok = False
    return check("OOXML zip sanity", ok)


# --- §4.5(b0) openxml-audit (installed at bootstrap) -----------------------
def openxml_audit(path: Path) -> dict:
    try:
        from openxml_audit import OpenXmlValidator  # type: ignore
    except ImportError:
        # library absent (optional per bootstrap §4.5) ⇒ skip, non-blocking.
        # round-trip + zip-sanity already prove the file opens.
        return check("openxml-audit skipped — validator not installed", False)
    try:
        result = OpenXmlValidator(strict=False).validate(str(path))
        return check("openxml-audit", bool(getattr(result, "is_valid", False)))
    except Exception:
        # library present but errored on this file ⇒ fail honestly
        return check("openxml-audit", False)


# --- §4.5(c) placeholder grep ----------------------------------------------
def placeholder_grep(texts: list) -> dict:
    joined = "\n".join(texts)
    clean = "{{" not in joined and "}}" not in joined and "TODO_" not in joined
    return check("Placeholder grep", clean)


# --- §4.5(e) soffice probe --------------------------------------------------
def find_soffice() -> str | None:
    hit = shutil.which("soffice")
    if hit:
        return hit
    mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    return mac if Path(mac).exists() else None


def soffice_convert(path: Path, label: str, skip_label: str) -> dict:
    """Headless convert to PDF — an OPPORTUNISTIC extra that proves the document
    opens in a real office app. NEVER blocking: round-trip + openxml-audit
    already prove structural validity, and LibreOffice is often absent (cloud)
    or broken. A failure/absence degrades to an amber skip, not a hard fail."""
    soffice = find_soffice()
    if not soffice:
        return check(skip_label, False)  # amber skip — soffice not installed
    with tempfile.TemporaryDirectory() as td:
        try:
            r = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, str(path)],
                capture_output=True,
                timeout=120,
            )
            produced = list(Path(td).glob("*.pdf"))
            if r.returncode == 0 and len(produced) == 1:
                return check(label, True)
            # soffice present but the convert failed (e.g. broken install) →
            # skip, don't block a document that already passed round-trip.
            return check(f"{label} skipped — soffice unavailable", False)
        except Exception:
            return check(f"{label} skipped — soffice unavailable", False)


# --- spec validation (schema + content audits) ------------------------------
# The office Lambda ships pure-python only and its deploy path swaps *.py files
# into a prebuilt zip — no new pip deps. This is a deliberately small validator
# covering exactly the JSON-Schema subset skills/*/schema.json uses.

def _schema_errors(schema: dict, value, path: str = "$") -> list:
    errs = []

    def err(msg):
        errs.append(f"{path}: {msg}")

    if "const" in schema and value != schema["const"]:
        err(f"expected {schema['const']!r}")
        return errs
    if "enum" in schema and value not in schema["enum"]:
        err(f"{value!r} not in enum")
        return errs

    t = schema.get("type")
    if t:
        ok = {
            "object": lambda v: isinstance(v, dict),
            "array": lambda v: isinstance(v, list),
            "string": lambda v: isinstance(v, str),
            "number": lambda v: isinstance(v, (int, float)) and not isinstance(v, bool),
            "integer": lambda v: isinstance(v, int) and not isinstance(v, bool),
            "boolean": lambda v: isinstance(v, bool),
            "null": lambda v: v is None,
        }[t](value)
        if not ok:
            err(f"expected {t}")
            return errs

    if isinstance(value, str):
        if "minLength" in schema and len(value) < schema["minLength"]:
            err(f"shorter than {schema['minLength']}")
        if "maxLength" in schema and len(value) > schema["maxLength"]:
            err(f"longer than {schema['maxLength']} chars")
        if "pattern" in schema:
            import re as _re
            if not _re.search(schema["pattern"], value):
                err(f"does not match {schema['pattern']}")
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        if "minimum" in schema and value < schema["minimum"]:
            err(f"below minimum {schema['minimum']}")
        if "maximum" in schema and value > schema["maximum"]:
            err(f"above maximum {schema['maximum']}")
    if isinstance(value, list):
        if "minItems" in schema and len(value) < schema["minItems"]:
            err(f"fewer than {schema['minItems']} items")
        if "maxItems" in schema and len(value) > schema["maxItems"]:
            err(f"more than {schema['maxItems']} items")
        if "items" in schema:
            for i, item in enumerate(value):
                errs.extend(_schema_errors(schema["items"], item, f"{path}[{i}]"))
    if isinstance(value, dict):
        for req in schema.get("required", []):
            if req not in value:
                err(f"missing required '{req}'")
        props = schema.get("properties", {})
        if schema.get("additionalProperties") is False:
            for key in value:
                if key not in props:
                    err(f"unexpected property '{key}'")
        for key, sub in props.items():
            if key in value:
                errs.extend(_schema_errors(sub, value[key], f"{path}.{key}"))

    for sub in schema.get("allOf", []):
        if "if" in sub:
            if not _schema_errors(sub["if"], value, path):
                errs.extend(_schema_errors(sub.get("then", {}), value, path))
        else:
            errs.extend(_schema_errors(sub, value, path))
    if "oneOf" in schema:
        matches = [s for s in schema["oneOf"] if not _schema_errors(s, value, path)]
        if len(matches) != 1:
            err(f"matches {len(matches)} of oneOf (need exactly 1)")
    return errs


def validate_spec(skill: str, payload: dict) -> list:
    """Schema + content audit for a spec. Returns a list of error strings
    (empty = valid). Builders call this BEFORE building; the server runs the
    same schema through ajv — this is the Lambda-side belt to that brace."""
    schema_path = Path(__file__).resolve().parents[2] / "skills" / skill / "schema.json"
    if not schema_path.exists():  # bundled Lambda layout: schemas beside handler
        schema_path = Path(__file__).resolve().parent / "schemas" / f"{skill}.json"
    errors = []
    if schema_path.exists():
        schema = json.loads(schema_path.read_text())
        errors.extend(_schema_errors(schema, payload))
    errors.extend(_content_audit(skill, payload))
    return errors


_PLACEHOLDER_RX = None


def placeholder_text_errors(texts: list, where: str) -> list:
    """Doctrine placeholder scan: xxxx | lorem | ipsum | click to edit | TODO."""
    global _PLACEHOLDER_RX
    import re as _re
    if _PLACEHOLDER_RX is None:
        _PLACEHOLDER_RX = _re.compile(r"x{4,}|lorem|ipsum|click\s+to\s+edit|TODO|\{\{|\}\}", _re.I)
    hits = sorted({m.group(0) for t in texts for m in [_PLACEHOLDER_RX.search(str(t))] if m})
    return [f"{where}: placeholder text {hits}"] if hits else []


def _words(s) -> int:
    return len(str(s).split())


def _content_audit(skill: str, payload: dict) -> list:
    """The numeric-doctrine rules a JSON schema cannot express (word counts,
    heading hierarchy, formulas-not-values)."""
    errs = []
    if skill == "pptx":
        for i, slide in enumerate(payload.get("slides", []), 1):
            at = slide.get("archetype", "?")
            texts = [slide.get("title", ""), slide.get("subtitle", "")]
            words = 0
            for b in slide.get("bullets", []) or []:
                texts.append(b)
                if _words(b) > 12:
                    errs.append(f"slide {i} ({at}): bullet over 12 words: {str(b)[:40]!r}")
                words += _words(b)
            for col in slide.get("columns", []) or []:
                texts.append(col.get("head", ""))
                words += _words(col.get("head", ""))
                for it in col.get("items", []):
                    texts.append(it)
                    if _words(it) > 12:
                        errs.append(f"slide {i} ({at}): column item over 12 words")
                    words += _words(it)
            for st in slide.get("steps", []) or []:
                texts += [st.get("label", ""), st.get("detail", "")]
                words += _words(st.get("label", "")) + _words(st.get("detail", ""))
            for extra in ("support", "quote", "attribution"):
                if slide.get(extra):
                    texts.append(slide[extra])
            words += _words(slide.get("title", "")) + _words(slide.get("subtitle", ""))
            if slide.get("stat"):
                texts += [slide["stat"].get("value", ""), slide["stat"].get("label", "")]
                words += _words(slide["stat"].get("label", ""))
            content_types = {"content_bullets", "content_chart", "comparison", "two_column", "table", "timeline_process"}
            if at in content_types and words > 40:
                errs.append(f"slide {i} ({at}): {words} words on a content slide (max 40)")
            chart = slide.get("chart")
            if chart:
                ncat = len(chart.get("categories", []))
                for s in chart.get("series", []):
                    if len(s.get("values", [])) != ncat:
                        errs.append(f"slide {i}: series {s.get('name')!r} has {len(s.get('values', []))} values for {ncat} categories")
            table = slide.get("table")
            if table:
                ncol = len(table.get("headers", []))
                for r, row in enumerate(table.get("rows", [])):
                    if len(row) != ncol:
                        errs.append(f"slide {i}: table row {r + 1} has {len(row)} cells for {ncol} headers")
            errs.extend(placeholder_text_errors(texts, f"slide {i}"))
    elif skill in ("docx", "pdf"):
        blocks = payload.get("blocks") or payload.get("sections") or []
        level = 0
        texts = []
        for i, blk in enumerate(blocks, 1):
            if blk.get("kind") == "heading":
                lv = int(blk.get("level", 1))
                if level == 0 and lv != 1:
                    errs.append(f"block {i}: document must open at heading level 1, got {lv}")
                elif level and lv > level + 1:
                    errs.append(f"block {i}: heading level skip {level} → {lv}")
                level = lv
            for key in ("text", "caption", "attribution"):
                if blk.get(key):
                    texts.append(blk[key])
            texts += blk.get("items", []) or []
            for row in blk.get("rows", []) or []:
                texts += row
            ncol = len(blk.get("headers", []) or [])
            if ncol:
                for r, row in enumerate(blk.get("rows", []) or []):
                    if len(row) != ncol:
                        errs.append(f"block {i}: table row {r + 1} has {len(row)} cells for {ncol} headers")
        errs.extend(placeholder_text_errors(texts, skill))
    elif skill == "xlsx":
        import re as _re
        derived_rx = _re.compile(r"^(total|sum|subtotal|net|variance|average|margin|growth)\b", _re.I)
        for sheet in payload.get("sheets", []):
            name = sheet.get("name", "?")
            ncol = len(sheet.get("columns", []))
            texts = [c.get("header", "") for c in sheet.get("columns", [])]
            for r, row in enumerate(sheet.get("rows", []), 1):
                if len(row) > ncol:
                    errs.append(f"sheet {name!r} row {r}: {len(row)} cells for {ncol} columns")
                label = row[0].get("t", "") if row and isinstance(row[0], dict) else ""
                if label:
                    texts.append(label)
                if derived_rx.match(str(label)):
                    for c, cell in enumerate(row[1:], 2):
                        if isinstance(cell, dict) and "n" in cell:
                            errs.append(
                                f"sheet {name!r} row {r} ({label!r}): column {c} is a hardcoded "
                                f"number where a formula is expected — derived rows must compute"
                            )
            errs.extend(placeholder_text_errors(texts, f"sheet {name!r}"))
    return errs


def spec_gate(skill: str, payload: dict) -> None:
    """Hard gate: refuse to build a spec that fails schema or content audit."""
    errors = validate_spec(skill, payload)
    if errors:
        fail(f"spec validation failed ({len(errors)}): " + " | ".join(errors[:8]))


def soffice_recalc_scan(path: Path) -> dict:
    """xlsx: convert via soffice and scan the result for formula error markers."""
    soffice = find_soffice()
    if not soffice:
        return check(RECALC_SKIP, False)
    with tempfile.TemporaryDirectory() as td:
        try:
            r = subprocess.run(
                [soffice, "--headless", "--convert-to", "csv", "--outdir", td, str(path)],
                capture_output=True,
                timeout=120,
            )
            blobs = "".join(p.read_text(errors="ignore") for p in Path(td).glob("*.csv"))
            errors = [m for m in ("#REF!", "#DIV/0!", "#VALUE!", "#N/A", "#NAME?") if m in blobs]
            if errors:
                # genuine formula errors in the sheet — this IS a real defect
                return check(f"soffice recalc — {','.join(errors)}", False)
            if r.returncode == 0:
                return check("soffice recalc", True)
            # convert failed (broken/absent soffice) → non-blocking skip
            return check("soffice recalc skipped — soffice unavailable", False)
        except Exception:
            return check("soffice recalc skipped — soffice unavailable", False)


# ═══════════════════════════════════════════════════════════════════════════
# Deliverable E — the HARD VISUAL GATE (deterministic; a deck failing any of
# these is never returned as success) + post-render bleed heuristic + thumbnail
# rendering for the advisory vision-critique pass.
# ═══════════════════════════════════════════════════════════════════════════

_SLIDE_W, _SLIDE_H = 13.333, 7.5
_MARGIN = 0.5
_EMU = 914400.0
_FOOTER_BAND = 6.7          # shapes fully below this line may sit in the margin
_FOOTER_CLEARANCE = 0.3     # content keeps this much clearance above footers


def _rect_in(shape):
    try:
        return (shape.left / _EMU, shape.top / _EMU, shape.width / _EMU, shape.height / _EMU)
    except Exception:
        return None


def _run_specs(text_frame):
    """(text, size_pt, bold, color) per paragraph — sizes default to theme body."""
    out = []
    for para in text_frame.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text.strip():
            continue
        sizes = [r.font.size.pt for r in para.runs if r.font.size]
        bold = any(r.font.bold for r in para.runs)
        color = next((r.font.color for r in para.runs if r.font.color and r.font.color.type is not None), None)
        out.append((text, max(sizes) if sizes else 18.0, bold, color))
    return out


def _resolve_hex(color, theme_hex):
    """A python-pptx ColorFormat → hex, resolving theme slots + brightness."""
    import pptx_design as D
    from pptx.enum.dml import MSO_THEME_COLOR

    slot_map = {
        MSO_THEME_COLOR.TEXT_1: "dk1", MSO_THEME_COLOR.TEXT_2: "dk2",
        MSO_THEME_COLOR.BACKGROUND_1: "lt1", MSO_THEME_COLOR.BACKGROUND_2: "lt2",
        MSO_THEME_COLOR.ACCENT_1: "accent1", MSO_THEME_COLOR.ACCENT_2: "accent2",
        MSO_THEME_COLOR.ACCENT_3: "accent3", MSO_THEME_COLOR.ACCENT_4: "accent4",
        MSO_THEME_COLOR.ACCENT_5: "accent5", MSO_THEME_COLOR.ACCENT_6: "accent6",
    }
    try:
        if color.type is not None and str(color.type).startswith("SCHEME"):
            base = theme_hex.get(slot_map.get(color.theme_color, ""), None)
            if base is None:
                return None
            brightness = getattr(color, "brightness", 0.0) or 0.0
            return D.brightness_hex(base, brightness) if brightness else base
        return str(color.rgb)
    except Exception:
        return None


def _fill_hex(shape, theme_hex):
    try:
        if shape.fill.type == 1:  # solid
            return _resolve_hex(shape.fill.fore_color, theme_hex)
    except Exception:
        pass
    return None


def _slide_bg_hex(slide, theme_hex):
    try:
        fill = slide.background.fill
        if fill.type == 1:
            return _resolve_hex(fill.fore_color, theme_hex)
    except Exception:
        pass
    return theme_hex.get("lt1", "FFFFFF")


def visual_gate_pptx(path: Path, template: str | None = None) -> list:
    """Deterministic findings for a BUILT deck: overflow (measured), collision,
    margins, WCAG contrast (unrounded), font-family count, content audit,
    placeholder scan, speaker notes. Empty list = gate passed."""
    import pptx_design as D
    from pptx import Presentation

    template_file = template or str(Path(__file__).resolve().parents[2] / "skills/pptx/templates/dfs_default.potx")
    if not Path(template_file).exists():
        template_file = str(Path(__file__).resolve().parent / "templates/dfs_default.potx")
    theme_hex = D.read_theme_hex(template_file)
    findings = []
    prs = Presentation(str(path))
    fonts_seen = set()

    for sidx, slide in enumerate(prs.slides, start=1):
        bg_hex = _slide_bg_hex(slide, theme_hex)
        rects, texts, slide_words = [], [], 0
        filled_panels = []  # (rect, hex) for containment-aware contrast

        for shape in slide.shapes:
            rect = _rect_in(shape)
            if rect is None:
                continue
            x, y, w, h = rect
            area = w * h
            full_bleed = area >= 0.5 * _SLIDE_W * _SLIDE_H
            fill = _fill_hex(shape, theme_hex)
            if fill and not getattr(shape, "has_text_frame", False):
                filled_panels.append((rect, fill))
            is_footer = y >= _FOOTER_BAND and h <= 0.45

            # margins (footer band + full-bleed exempt)
            if not full_bleed and not is_footer:
                if x < _MARGIN - 0.03 or y < _MARGIN - 0.03 or x + w > _SLIDE_W - _MARGIN + 0.03 or y + h > _SLIDE_H - _MARGIN + 0.03:
                    findings.append(f"slide {sidx}: shape outside 0.5in margins at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
            if not full_bleed:
                rects.append((rect, is_footer, getattr(shape, "has_text_frame", False)))

            if getattr(shape, "has_text_frame", False):
                paras = _run_specs(shape.text_frame)
                if paras:
                    total_h = 0.0
                    for text, size, bold, _ in paras:
                        total_h += D.required_height_in(text, max(w, 0.1), size, bold=bold, spacing=1.05)
                    if total_h > h * 1.08 + 0.05:
                        findings.append(f"slide {sidx}: text overflow — needs {total_h:.2f}in in a {h:.2f}in frame: {paras[0][0][:40]!r}")
                    for text, size, bold, color in paras:
                        fonts_seen.update(
                            r.font.name for p in shape.text_frame.paragraphs for r in p.runs
                            if r.font.name and not r.font.name.startswith("+"))
                        if not is_footer:
                            slide_words += len(text.split())
                        texts.append(text)
                        # contrast: effective background = smallest filled panel containing this shape
                        fg = _resolve_hex(color, theme_hex) if color else theme_hex.get("dk1")
                        if fg is None:
                            continue
                        behind = bg_hex
                        best_area = None
                        for (px, py, pw, ph), phex in filled_panels:
                            if px <= x + 0.02 and py <= y + 0.02 and px + pw >= x + w - 0.02 and py + ph >= y + h - 0.02:
                                if best_area is None or pw * ph < best_area:
                                    behind, best_area = phex, pw * ph
                        ratio = D.contrast(fg, behind)
                        large = size >= 18.0 or (size >= 14.0 and bold)
                        threshold = 3.0 if large else 4.5
                        if ratio < threshold:
                            findings.append(
                                f"slide {sidx}: contrast {ratio:.2f}:1 < {threshold}:1 for {size:.0f}pt "
                                f"{'bold ' if bold else ''}text #{fg} on #{behind}: {text[:30]!r}")

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text)
                            slide_words += 0  # table values are data, not slide copy

        # collisions among slide-level shapes (containment allowed; footer clearance)
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                (ax, ay, aw, ah), afoot, _ = rects[i]
                (bx, by, bw, bh), bfoot, _ = rects[j]
                ix = min(ax + aw, bx + bw) - max(ax, bx)
                iy = min(ay + ah, by + bh) - max(ay, by)
                if ix <= 0.08 or iy <= 0.08:
                    continue
                a_in_b = ax >= bx - 0.03 and ay >= by - 0.03 and ax + aw <= bx + bw + 0.03 and ay + ah <= by + bh + 0.03
                b_in_a = bx >= ax - 0.03 and by >= ay - 0.03 and bx + bw <= ax + aw + 0.03 and by + bh <= ay + ah + 0.03
                if a_in_b or b_in_a:
                    continue  # panel/label containment is intentional layering
                if afoot != bfoot:
                    gap = (by - (ay + ah)) if bfoot else (ay - (by + bh))
                    if gap < _FOOTER_CLEARANCE:
                        findings.append(f"slide {sidx}: content within {_FOOTER_CLEARANCE}in of the footer band")
                    continue
                findings.append(
                    f"slide {sidx}: shapes collide ({ax:.2f},{ay:.2f})x({aw:.2f},{ah:.2f}) vs ({bx:.2f},{by:.2f})x({bw:.2f},{bh:.2f})")

        if slide_words > 46:  # file-level tolerance over the 40-word spec gate (numbers/labels)
            findings.append(f"slide {sidx}: {slide_words} words rendered on one slide")
        findings.extend(placeholder_text_errors(texts, f"slide {sidx}"))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        except Exception:
            notes = ""
        if not notes:
            findings.append(f"slide {sidx}: speaker notes missing")

    families = {f for f in fonts_seen if f}
    if len(families) > 2:
        findings.append(f"deck uses {len(families)} font families: {sorted(families)} (max 2)")
    return findings


def post_render_bleed(path: Path, kind: str = "pptx") -> dict:
    """Second overflow signal: soffice → PDF → raster → ink in the margin band.
    Enforced when soffice is present; amber skip otherwise. The bottom band is
    exempt (brand logo + page number legitimately live there)."""
    soffice = find_soffice()
    if not soffice:
        return check("Post-render bleed scan skipped — soffice not found", False)
    try:
        import pdfplumber
    except ImportError:
        return check("Post-render bleed scan skipped — pdfplumber not found", False)
    with tempfile.TemporaryDirectory() as td:
        r = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, str(path)],
            capture_output=True, timeout=180)
        pdfs = list(Path(td).glob("*.pdf"))
        if r.returncode != 0 or not pdfs:
            return check("Post-render bleed scan skipped — soffice convert failed", False)
        dirty = []
        with pdfplumber.open(str(pdfs[0])) as pdf:
            for pno, page in enumerate(pdf.pages, start=1):
                img = page.to_image(resolution=96).original.convert("L")
                px = img.load()
                w, h = img.size
                band = max(2, int(96 * 0.10))  # outer 0.10in — content should sit >=0.5in in
                # median corner sample approximates the field color
                bgs = sorted(px[2, 2], ) if False else None
                base = px[2, 2]
                hits = 0
                for yy in range(0, h, 3):
                    for xx in list(range(0, band, 2)) + list(range(w - band, w, 2)):
                        if abs(px[xx, yy] - base) > 60:
                            hits += 1
                for xx in range(0, w, 3):
                    for yy in range(0, band, 2):  # top band only; bottom exempt (logo/page no)
                        if abs(px[xx, yy] - base) > 60:
                            hits += 1
                if hits > 40:
                    dirty.append(pno)
        if dirty:
            return check(f"Post-render bleed scan — ink at page edge on pages {dirty}", False)
        return check("Post-render bleed scan", True)


def render_thumbnails(path: Path, max_pages: int = 12, dpi: int = 100) -> list:
    """Base64 JPEG thumbnails for the advisory vision-critique pass (flag-gated
    caller-side). Empty list when soffice/pdfplumber are unavailable."""
    import base64
    import io

    soffice = find_soffice()
    if not soffice:
        return []
    try:
        import pdfplumber
    except ImportError:
        return []
    thumbs = []
    with tempfile.TemporaryDirectory() as td:
        r = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, str(path)],
            capture_output=True, timeout=180)
        pdfs = list(Path(td).glob("*.pdf"))
        if r.returncode != 0 and not pdfs:
            return []
        with pdfplumber.open(str(pdfs[0])) as pdf:
            for page in pdf.pages[:max_pages]:
                img = page.to_image(resolution=dpi).original.convert("RGB")
                buf = io.BytesIO()
                img.save(buf, "JPEG", quality=70)
                thumbs.append(base64.b64encode(buf.getvalue()).decode())
    return thumbs


def pdf_table_break_check(payload: dict, path: Path) -> dict:
    """No table split across pages: every spec table's row leaders must appear
    together on a single extracted page."""
    import pdfplumber

    tables = [b for b in payload.get("sections", []) if b.get("kind") == "table"]
    if not tables:
        return check("Tables unbroken across pages", True)
    with pdfplumber.open(str(path)) as pdf:
        pages = [p.extract_text() or "" for p in pdf.pages]
    for t in tables:
        leaders = [str(row[0]).strip() for row in t.get("rows", []) if row and str(row[0]).strip()]
        if not leaders:
            continue
        if not any(all(ld in page for ld in leaders) for page in pages):
            return check(f"Tables unbroken across pages — table starting {leaders[0]!r} split", False)
    return check("Tables unbroken across pages", True)

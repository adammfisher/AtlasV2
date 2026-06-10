#!/usr/bin/env python3
"""Build the default Atlas office templates (PRD §5). Idempotent.

Theme matches client/src/theme/tokens.ts: accent coral #d97757 on warm charcoal.
The .potx/.dotx files are structurally pptx/docx packages (python-pptx/docx do
not write template content-types); helpers open them by path, so this is
transparent to the pipeline.
"""
from pathlib import Path

ACCENT = "D97757"
CHARCOAL = "262624"
TEXT_DARK = "1B1A18"
REPO = Path(__file__).resolve().parents[2]

def make_potx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()  # default master has title/bullets/two-content/blank layouts
    # Stamp the theme accent into the master title placeholder + a footer bar so
    # decks visibly carry the Atlas theme even before per-slide styling.
    master = prs.slide_masters[0]
    for ph in master.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            for para in ph.text_frame.paragraphs:
                para.font.color.rgb = RGBColor.from_string(ACCENT)
                para.font.name = "Helvetica Neue"
    # stamp the accent into the theme part (accent1) so charts/shapes inherit it
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    theme_part = master.part.part_related_by(RT.THEME)
    blob = theme_part.blob
    # default Office theme accent1 is 4472C4; swap for Atlas coral
    theme_part._blob = blob.replace(b'val="4472C4"', f'val="{ACCENT}"'.encode())
    prs.save(str(path))


def make_dotx(path: Path) -> None:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    styles = doc.styles
    for name, size, bold, color in [
        ("Heading 1", 20, True, ACCENT),
        ("Heading 2", 15, True, TEXT_DARK),
        ("Heading 3", 12, True, TEXT_DARK),
        ("Normal", 10.5, False, TEXT_DARK),
    ]:
        st = styles[name]
        st.font.size = Pt(size)
        st.font.bold = bold
        st.font.name = "Helvetica Neue"
        st.font.color.rgb = RGBColor.from_string(color)
    doc.save(str(path))


def make_xlsx_theme(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, NamedStyle

    wb = Workbook()
    header = NamedStyle(name="atlas_header")
    header.font = Font(bold=True, color="FFFFFF", name="Helvetica Neue")
    header.fill = PatternFill("solid", fgColor=ACCENT)
    wb.add_named_style(header)
    body = NamedStyle(name="atlas_body")
    body.font = Font(color=TEXT_DARK, name="Helvetica Neue")
    wb.add_named_style(body)
    ws = wb.active
    ws.title = "Theme"
    ws["A1"] = "Atlas default theme"
    ws["A1"].style = "atlas_header"
    wb.save(str(path))


def main() -> None:
    targets = {
        REPO / "skills/pptx/templates/atlas_default.potx": make_potx,
        REPO / "skills/docx/templates/atlas_default.dotx": make_dotx,
        REPO / "skills/xlsx/templates/atlas_default.xlsx": make_xlsx_theme,
    }
    for path, builder in targets.items():
        path.parent.mkdir(parents=True, exist_ok=True)
        if path.exists():
            print(f"exists, skipping: {path.relative_to(REPO)}")
            continue
        builder(path)
        print(f"built: {path.relative_to(REPO)}")


if __name__ == "__main__":
    main()

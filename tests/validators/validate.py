#!/usr/bin/env python3
"""File-validity harness (TESTPLAN §4).

Usage: validate.py <kind> <file> [--spec spec.json]
Prints one JSON line: {"ok": bool, "kind": str, "findings": [str, ...]}
Exit 0 when ok, 1 when findings, 2 on harness error.

kinds: pptx docx xlsx pdf md svg  (native)
       mermaid react site         (delegated to js-validate.ts, which reuses
                                   the product's own validators + esbuild)

spec.json (all keys optional):
  {"slides": 5, "contains": ["text", ...], "sheets": ["Budget"], "columns": ["Category"]}
"""
import argparse
import json
import re
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "scripts" / "office"))

PLACEHOLDER_RE = re.compile(r"x{4,}|lorem|ipsum|click to edit|TODO|\{\{|\}\}", re.IGNORECASE)
XLSX_ERRORS = ("#REF!", "#DIV/0!", "#NAME?", "#VALUE!", "#N/A")


def v_pptx(path: Path, spec: dict) -> list[str]:
    from pptx import Presentation  # type: ignore

    findings: list[str] = []
    try:
        prs = Presentation(str(path))
    except Exception as e:  # noqa: BLE001
        return [f"cannot open pptx: {e}"]
    n = len(prs.slides)
    if n == 0:
        findings.append("zero slides")
    if "slides" in spec and n != spec["slides"]:
        findings.append(f"slide count {n} != requested {spec['slides']}")
    all_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            all_text.append(text)
            if PLACEHOLDER_RE.search(text):
                findings.append(f"slide {i}: leftover placeholder text: {text[:60]!r}")
            if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0:
                title = text
        if i > 1 and not title.strip():
            findings.append(f"slide {i}: empty title placeholder")
    joined = "\n".join(all_text)
    for want in spec.get("contains", []):
        if want.lower() not in joined.lower():
            findings.append(f"missing requested content: {want!r}")
    if spec.get("_design"):
        # the product's own deterministic visual gate (AX-design): overflow,
        # collisions, margins, WCAG contrast, font families, placeholder scan
        try:
            import validate_common as vc  # type: ignore

            findings.extend(f"design: {f}" for f in vc.visual_gate_pptx(path))
        except Exception as e:  # noqa: BLE001
            findings.append(f"design gate could not run: {e}")
    return findings


def v_docx(path: Path, spec: dict) -> list[str]:
    import docx  # type: ignore

    findings: list[str] = []
    try:
        d = docx.Document(str(path))
    except Exception as e:  # noqa: BLE001
        return [f"cannot open docx: {e}"]
    styles = [p.style.name or "" for p in d.paragraphs]
    if not any(s.startswith(("Heading", "Title")) for s in styles):
        findings.append("no heading structure")
    text = "\n".join(p.text for p in d.paragraphs)
    for want in spec.get("contains", []):
        if want.lower() not in text.lower():
            findings.append(f"missing requested content: {want!r}")
    return findings


def v_xlsx(path: Path, spec: dict) -> list[str]:
    import openpyxl  # type: ignore

    findings: list[str] = []
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True)
    except Exception as e:  # noqa: BLE001
        return [f"cannot open xlsx: {e}"]
    # recalculated values: prefer the product's soffice recalc scan when available.
    # It returns a single check dict {label, ok}; a "skipped" (no soffice) result
    # is NOT a defect — fall back to scanning cached values instead.
    recalc_done = False
    try:
        import validate_common as vc  # type: ignore

        res = vc.soffice_recalc_scan(path)
        label = str(res.get("label", ""))
        if res.get("ok"):
            recalc_done = True
        elif "skip" not in label.lower() and "unavailable" not in label.lower():
            findings.append(f"recalc: {label}")
            recalc_done = True
    except Exception:
        recalc_done = False
    if not recalc_done:
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and any(t in cell.value for t in XLSX_ERRORS):
                        findings.append(f"{ws.title}!{cell.coordinate}: formula error {cell.value!r}")
    names = [ws.title for ws in wb.worksheets]
    for want in spec.get("sheets", []):
        if want not in names:
            findings.append(f"missing sheet {want!r}")
    if spec.get("columns"):
        headers: set[str] = set()
        for ws in wb.worksheets:
            for row in ws.iter_rows(max_row=1):
                headers.update(str(c.value) for c in row if c.value is not None)
        for want in spec["columns"]:
            if want not in headers:
                findings.append(f"missing column header {want!r}")
    return findings


def v_pdf(path: Path, spec: dict) -> list[str]:
    import pdfplumber  # type: ignore

    findings: list[str] = []
    try:
        with pdfplumber.open(str(path)) as pdf:
            if len(pdf.pages) == 0:
                findings.append("zero pages")
            text = "\n".join((p.extract_text() or "") for p in pdf.pages)
    except Exception as e:  # noqa: BLE001
        return [f"cannot open pdf: {e}"]
    if not text.strip():
        findings.append("no extractable text")
    for want in spec.get("contains", []):
        if want.lower() not in text.lower():
            findings.append(f"missing requested content: {want!r}")
    return findings


def v_md(path: Path, spec: dict) -> list[str]:
    findings: list[str] = []
    text = path.read_text(encoding="utf-8", errors="replace")
    if text.count("```") % 2 != 0:
        findings.append("unclosed code fence")
    if not re.search(r"^#{1,6} ", text, re.MULTILINE):
        findings.append("no headings")
    for want in spec.get("contains", []):
        if want.lower() not in text.lower():
            findings.append(f"missing requested content: {want!r}")
    return findings


def v_svg(path: Path, spec: dict) -> list[str]:
    import xml.etree.ElementTree as ET

    findings: list[str] = []
    try:
        root = ET.fromstring(path.read_text(encoding="utf-8"))
    except ET.ParseError as e:
        return [f"invalid XML: {e}"]
    if not root.tag.endswith("svg"):
        findings.append(f"root element is {root.tag!r}, not svg")
    viewbox = root.get("viewBox")
    if not viewbox:
        findings.append("missing viewBox")
    else:
        parts = viewbox.replace(",", " ").split()
        if len(parts) != 4 or float(parts[2]) <= 0 or float(parts[3]) <= 0:
            findings.append(f"degenerate viewBox {viewbox!r}")
    return findings


def v_delegated(kind: str, path: Path, spec: dict) -> list[str]:
    """mermaid / react / site — product validators via node (js-validate.ts)."""
    args = [str(ROOT / "node_modules" / ".bin" / "tsx"), str(ROOT / "tests" / "validators" / "js-validate.ts"), kind, str(path)]
    if spec.get("contains"):
        import tempfile

        with tempfile.NamedTemporaryFile("w", suffix=".json", delete=False) as f:
            json.dump({"contains": spec["contains"]}, f)
            args.append(f.name)
    proc = subprocess.run(
        args,
        capture_output=True,
        text=True,
        timeout=120,
        cwd=str(ROOT),
    )
    if proc.returncode not in (0, 1):
        return [f"js-validate crashed: {proc.stderr.strip()[-300:]}"]
    try:
        verdict = json.loads(proc.stdout.strip().splitlines()[-1])
    except Exception:  # noqa: BLE001
        return [f"js-validate emitted no verdict: {proc.stdout[-200:]!r}"]
    return list(verdict.get("findings", []))


VALIDATORS = {"pptx": v_pptx, "docx": v_docx, "xlsx": v_xlsx, "pdf": v_pdf, "md": v_md, "svg": v_svg}


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("kind")
    ap.add_argument("file")
    ap.add_argument("--spec", default=None)
    ap.add_argument("--design", action="store_true", help="pptx: also run the product's deterministic visual gate")
    args = ap.parse_args()
    spec = json.loads(Path(args.spec).read_text()) if args.spec else {}
    if args.design:
        spec["_design"] = True
    path = Path(args.file)
    if not path.exists():
        print(json.dumps({"ok": False, "kind": args.kind, "findings": ["file does not exist"]}))
        return 1
    try:
        if args.kind in VALIDATORS:
            findings = VALIDATORS[args.kind](path, spec)
        elif args.kind in ("mermaid", "react", "site"):
            findings = v_delegated(args.kind, path, spec)
        else:
            print(json.dumps({"ok": False, "kind": args.kind, "findings": [f"unknown kind {args.kind!r}"]}))
            return 2
    except Exception as e:  # noqa: BLE001
        print(json.dumps({"ok": False, "kind": args.kind, "findings": [f"harness error: {e}"]}))
        return 2
    print(json.dumps({"ok": not findings, "kind": args.kind, "findings": findings}))
    return 0 if not findings else 1


if __name__ == "__main__":
    sys.exit(main())

#!/usr/bin/env python3
"""Phase 1 gate: the validity harness must accept a known-good file of every
kind and reject a deliberately corrupted one. Also (re)generates the seed
upload fixtures under tests/fixtures/files/.

Run: runtimes/python/venv/bin/python tests/validators/selftest.py
Exit 0 = all green.
"""
import json
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
PY = sys.executable
VALIDATE = ROOT / "tests" / "validators" / "validate.py"
FILES = ROOT / "tests" / "fixtures" / "files"
FILES.mkdir(parents=True, exist_ok=True)


def run_validate(kind: str, path: Path) -> dict:
    proc = subprocess.run([PY, str(VALIDATE), kind, str(path)], capture_output=True, text=True, timeout=180)
    try:
        return json.loads(proc.stdout.strip().splitlines()[-1])
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "findings": [f"no verdict: {e}; stdout={proc.stdout[-200:]!r} stderr={proc.stderr[-200:]!r}"]}


# ---------- known-good generators (double as seed upload fixtures) ----------

def good_pptx(out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Quarterly revenue grew 12% on enterprise expansion"
    s1.placeholders[1].text = "Sample deck for harness self-test"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Three drivers explain the growth"
    s2.placeholders[1].text_frame.text = "Enterprise upsell led bookings"
    s2.placeholders[1].text_frame.add_paragraph().text = "Churn fell to 2.1%"
    prs.save(str(out))


def good_docx(out: Path) -> None:
    import docx

    d = docx.Document()
    d.add_heading("Migration Runbook", level=1)
    d.add_paragraph("This runbook covers the six-week data migration engagement.")
    d.add_heading("Scope", level=2)
    d.add_paragraph("Extract, transform, and load the legacy ledger.")
    d.save(str(out))


def good_xlsx(out: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Category", "Plan", "Actual", "Variance"])
    ws.append(["Travel", 1000, 800, "=C2-B2"])
    ws.append(["Software", 2000, 2400, "=C3-B3"])
    ws.append(["Total", "=SUM(B2:B3)", "=SUM(C2:C3)", "=C4-B4"])
    wb.save(str(out))


def good_pdf(out: Path) -> None:
    from weasyprint import HTML

    HTML(string="<h1>Onboarding Checklist</h1><p>Welcome to the Axiom pilot program.</p>").write_pdf(str(out))


def good_md(out: Path) -> None:
    out.write_text("# Notes\n\nSome content.\n\n```py\nprint('hi')\n```\n", encoding="utf-8")


def good_svg(out: Path) -> None:
    out.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><circle cx="50" cy="50" r="40" fill="#26A697"/></svg>',
        encoding="utf-8",
    )


def good_mermaid(out: Path) -> None:
    out.write_text("flowchart TD\n  A[Start] --> B[Finish]\n", encoding="utf-8")


def good_react(out: Path) -> None:
    out.write_text(
        json.dumps({"entry": "/App.jsx", "files": {"/App.jsx": "export default function App() {\n  return <div>ok</div>;\n}\n"}}),
        encoding="utf-8",
    )


def good_site(out: Path) -> None:
    out.write_text(
        json.dumps({"files": {"/index.html": "<!doctype html>\n<html><body><h1>Site</h1></body></html>\n"}}),
        encoding="utf-8",
    )


# ---------- corrupted generators ----------

def truncate(src: Path, out: Path) -> None:
    data = src.read_bytes()
    out.write_bytes(data[: max(10, len(data) // 3)])


def bad_pptx(tmp: Path) -> Path:
    good = tmp / "g.pptx"
    good_pptx(good)
    bad = tmp / "bad.pptx"
    truncate(good, bad)
    return bad


def bad_pptx_semantic(tmp: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Lorem ipsum click to edit"
    s.placeholders[1].text_frame.text = "TODO: fill this in"
    out = tmp / "bad-semantic.pptx"
    prs.save(str(out))
    return out


def bad_docx(tmp: Path) -> Path:
    good = tmp / "g.docx"
    good_docx(good)
    bad = tmp / "bad.docx"
    truncate(good, bad)
    return bad


def bad_xlsx(tmp: Path) -> Path:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Header"])
    ws.append(["#REF!"])
    out = tmp / "bad.xlsx"
    wb.save(str(out))
    return out


def bad_pdf(tmp: Path) -> Path:
    good = tmp / "g.pdf"
    good_pdf(good)
    bad = tmp / "bad.pdf"
    truncate(good, bad)
    return bad


def bad_md(tmp: Path) -> Path:
    out = tmp / "bad.md"
    out.write_text("# Title\n\n```py\nprint('unclosed fence'\n", encoding="utf-8")
    return out


def bad_svg(tmp: Path) -> Path:
    out = tmp / "bad.svg"
    out.write_text('<svg xmlns="http://www.w3.org/2000/svg"><circle cx="50"', encoding="utf-8")
    return out


def bad_mermaid(tmp: Path) -> Path:
    out = tmp / "bad.mmd"
    out.write_text("this is not a diagram type\nA --> B\n", encoding="utf-8")
    return out


def bad_react(tmp: Path) -> Path:
    out = tmp / "bad-react.json"
    out.write_text(json.dumps({"entry": "/App.jsx", "files": {"/App.jsx": "export default function App( { return <div>; }"}}), encoding="utf-8")
    return out


def bad_site(tmp: Path) -> Path:
    out = tmp / "bad-site.json"
    out.write_text(json.dumps({"files": {"/about.html": "<p>no index</p>"}}), encoding="utf-8")
    return out


GOOD = {
    "pptx": good_pptx, "docx": good_docx, "xlsx": good_xlsx, "pdf": good_pdf,
    "md": good_md, "svg": good_svg, "mermaid": good_mermaid, "react": good_react, "site": good_site,
}
BAD = {
    "pptx": bad_pptx, "docx": bad_docx, "xlsx": bad_xlsx, "pdf": bad_pdf,
    "md": bad_md, "svg": bad_svg, "mermaid": bad_mermaid, "react": bad_react, "site": bad_site,
}
SEED_EXT = {"pptx": "sample.pptx", "docx": "sample.docx", "xlsx": "sample.xlsx", "pdf": "sample.pdf"}


def make_seed_extras() -> None:
    (FILES / "sample.csv").write_text(
        "region,revenue\nEMEA,1200\nAMER,3400\nAPAC,900\n", encoding="utf-8"
    )
    from PIL import Image

    img = Image.new("RGB", (64, 64), (38, 166, 151))
    img.save(FILES / "sample.png")


def main() -> int:
    failures: list[str] = []
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        for kind, gen in GOOD.items():
            seed_name = SEED_EXT.get(kind)
            ext = {"mermaid": "mmd", "react": "json", "site": "json"}.get(kind, kind)
            target = FILES / seed_name if seed_name else tmp / f"good.{ext}"
            gen(target)
            verdict = run_validate(kind, target)
            status = "PASS" if verdict.get("ok") else "FAIL"
            print(f"  good {kind:8s} → {status}  {verdict.get('findings') or ''}")
            if not verdict.get("ok"):
                failures.append(f"good {kind} rejected: {verdict.get('findings')}")

        for kind, gen in BAD.items():
            bad = gen(tmp)
            verdict = run_validate(kind, bad)
            status = "PASS" if not verdict.get("ok") else "FAIL"
            print(f"  bad  {kind:8s} → {status}  {verdict.get('findings') or '(accepted!)'}")
            if verdict.get("ok"):
                failures.append(f"corrupted {kind} accepted")

        semantic = bad_pptx_semantic(tmp)
        verdict = run_validate("pptx", semantic)
        print(f"  bad  pptx-sem → {'PASS' if not verdict.get('ok') else 'FAIL'}  {verdict.get('findings')}")
        if verdict.get("ok"):
            failures.append("semantic-bad pptx (placeholder text) accepted")

    make_seed_extras()
    print(f"\nseed fixtures written to {FILES}")
    if failures:
        print("VALIDATOR SELFTEST: FAIL")
        for f in failures:
            print(f"  - {f}")
        return 1
    print("VALIDATOR SELFTEST: PASS (9 kinds accept-good + reject-bad, plus semantic pptx)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
